"""Genera PDFs y presentaciones por clase para el currículo v3.

Qué resuelve:
    El script viejo `generate_class_assets.py` asume estructura v1 con
    múltiples archivos por clase (`slides.md`, `teoria.md`, `ejercicios.md`,
    `homework.md`). El currículo v3 (232 clases, 9 partes) consolida toda
    la pauta en SOLO dos archivos: `README.md` (110 líneas, secciones H2
    estándar) y `notebook.ipynb` (12-16 celdas nbformat).

    Este script `_v3` lee esa nueva estructura y produce, por cada clase:

        1. `clase-<class_basename>-guia-explicativa.pdf` — guía explicativa
           generada con `scripts.generar_pdf_documento.render_markdown_text`.
        2. `clase-<class_basename>-presentacion.pptx` — deck de 7-8 slides
           con un layout coherente (portada, objetivo + outcomes, temas,
           definiciones, código de muestra, ejercicios + homework,
           referencias, cierre).

    Cada archivo se escribe dos veces: dentro de la propia carpeta de clase
    (para `content_loader.resolve_class_asset_path`) y en el mirror público
    `docs/pdfs/classes/` / `docs/presentaciones/classes/`.

Uso:
    python scripts/generate_class_assets_v3.py
    python scripts/generate_class_assets_v3.py --parte 0
    python scripts/generate_class_assets_v3.py --clase parte-0-prerrequisitos/006-python-tipos-estructuras-control-de-flujo
    python scripts/generate_class_assets_v3.py --force
"""

from __future__ import annotations

import argparse
import json
import re
import sys
import traceback
import unicodedata
from pathlib import Path

# Importar el renderer PDF ya existente del repo.
sys.path.insert(0, str(Path(__file__).resolve().parent))
from generar_pdf_documento import render_markdown_text  # noqa: E402
from pptx import Presentation  # noqa: E402
from pptx.dml.color import RGBColor  # noqa: E402
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE  # noqa: E402
from pptx.enum.text import MSO_ANCHOR  # noqa: E402
from pptx.util import Inches, Pt  # noqa: E402

try:
    import nbformat  # type: ignore
except ImportError:  # pragma: no cover - degradación
    nbformat = None  # type: ignore

BASE_DIR = Path(__file__).resolve().parents[1]
CLASSES_DIR = BASE_DIR / "classes"
PDF_OUTPUT_DIR = BASE_DIR / "docs" / "pdfs" / "classes"
PPTX_OUTPUT_DIR = BASE_DIR / "docs" / "presentaciones" / "classes"

# Paleta consistente con el deck v1 (re-usada para mantener identidad visual).
COLOR_BG = RGBColor(248, 250, 252)
COLOR_PANEL = RGBColor(255, 255, 255)
COLOR_PANEL_BORDER = RGBColor(203, 213, 225)
COLOR_HEADER_BAR = RGBColor(15, 23, 42)
COLOR_ACCENT = RGBColor(34, 197, 94)
COLOR_TEXT = RGBColor(15, 23, 42)
COLOR_MUTED = RGBColor(71, 85, 105)
COLOR_CODE = RGBColor(15, 23, 42)
COLOR_CODE_TEXT = RGBColor(248, 250, 252)

EMOJI_RE = re.compile(r"[\U0001F300-\U0001FAFF☀-➿⌀-⏿]")

PARTE_DIR_RE = re.compile(r"^parte-(\d+)-")
CLASS_DIR_RE = re.compile(r"^(\d{3})-")


# ---------------------------------------------------------------------------
# Helpers de texto / markdown
# ---------------------------------------------------------------------------


def read_text(path: Path) -> str:
    """Lee un archivo de texto en UTF-8."""
    return path.read_text(encoding="utf-8")


def clean_text(text: str) -> str:
    """Limpia markdown inline simple y emojis para reutilizar texto en PPTX."""
    clean = text
    clean = re.sub(r"\[(.+?)\]\(.+?\)", r"\1", clean)
    clean = re.sub(r"`(.+?)`", r"\1", clean)
    clean = re.sub(r"\*\*(.+?)\*\*", r"\1", clean)
    clean = re.sub(r"\*(.+?)\*", r"\1", clean)
    clean = EMOJI_RE.sub("", clean)
    clean = clean.replace("‍", "")
    clean = clean.replace("️", "")
    return clean.strip()


def normalize_heading(text: str) -> str:
    """Normaliza un heading para comparaciones tolerantes a emojis/acentos."""
    plain = clean_text(text).lower()
    plain = unicodedata.normalize("NFD", plain)
    plain = "".join(char for char in plain if unicodedata.category(char) != "Mn")
    return re.sub(r"\s+", " ", plain).strip()


# ---------------------------------------------------------------------------
# Descubrimiento de clases v3
# ---------------------------------------------------------------------------


def discover_class_dirs() -> list[Path]:
    """Devuelve todas las carpetas de clase v3 ordenadas por (parte, clase)."""
    dirs: list[tuple[int, int, Path]] = []
    for parte_dir in sorted(CLASSES_DIR.iterdir()):
        if not parte_dir.is_dir():
            continue
        m_parte = PARTE_DIR_RE.match(parte_dir.name)
        if not m_parte:
            continue
        parte_num = int(m_parte.group(1))
        for class_dir in sorted(parte_dir.iterdir()):
            if not class_dir.is_dir():
                continue
            m_class = CLASS_DIR_RE.match(class_dir.name)
            if not m_class:
                continue
            if not (class_dir / "README.md").exists():
                continue
            dirs.append((parte_num, int(m_class.group(1)), class_dir))
    dirs.sort(key=lambda item: (item[0], item[1]))
    return [item[2] for item in dirs]


def filter_dirs(
    dirs: list[Path],
    parte: int | None,
    clase: str | None,
) -> list[Path]:
    """Filtra el listado por --parte y/o --clase según CLI."""
    result = dirs
    if parte is not None:
        prefix = f"parte-{parte}-"
        result = [d for d in result if d.parent.name.startswith(prefix)]
    if clase:
        # Slug puede llegar como "parte-0-.../006-...-control-de-flujo" o con backslash.
        slug = clase.replace("\\", "/").strip("/")
        result = [
            d
            for d in result
            if d.name == slug
            or f"{d.parent.name}/{d.name}" == slug
            or d.name.startswith(slug)
        ]
    return result


# ---------------------------------------------------------------------------
# Parser de README.md v3
# ---------------------------------------------------------------------------


def parse_class_readme(readme_path: Path) -> dict:
    """Parsea un README.md v3 a un dict con título, subtítulo y secciones H2."""
    text = read_text(readme_path)
    lines = text.splitlines()

    title = readme_path.parent.name
    subtitle_parts: list[str] = []
    body_start = 0

    # Localizar H1.
    for idx, line in enumerate(lines):
        if line.startswith("# "):
            title = clean_text(line[2:])
            body_start = idx + 1
            break

    # Recolectar líneas blockquote inmediatas como subtítulo.
    cursor = body_start
    while cursor < len(lines):
        stripped = lines[cursor].strip()
        if stripped.startswith(">"):
            subtitle_parts.append(clean_text(stripped.lstrip(">").strip()))
            cursor += 1
        elif stripped == "":
            cursor += 1
        elif stripped.startswith("---"):
            cursor += 1
        else:
            break

    subtitle = " · ".join(part for part in subtitle_parts if part)

    # Partir en secciones H2.
    sections: list[tuple[str, str]] = []
    current_h2: str | None = None
    current_lines: list[str] = []

    for line in lines[cursor:]:
        if line.startswith("## "):
            if current_h2 is not None:
                sections.append((current_h2, "\n".join(current_lines).strip()))
            current_h2 = line[3:].strip()
            current_lines = []
        else:
            current_lines.append(line)

    if current_h2 is not None:
        sections.append((current_h2, "\n".join(current_lines).strip()))

    return {"title": title, "subtitle": subtitle, "sections": sections}


def get_section(readme_data: dict, *needles: str) -> str:
    """Devuelve el body de la primera sección cuyo H2 contiene cualquier needle."""
    targets = [normalize_heading(n) for n in needles]
    for h2, body in readme_data["sections"]:
        norm = normalize_heading(h2)
        if any(t in norm for t in targets):
            return body
    return ""


def list_items(markdown_text: str) -> list[str]:
    """Extrae bullets y listas numeradas — para Resultados, Ejercicios, etc."""
    items: list[str] = []
    for line in markdown_text.splitlines():
        stripped = line.strip()
        if stripped.startswith("- ") or stripped.startswith("* "):
            items.append(clean_text(stripped[2:]))
        elif re.match(r"^\d+\.\s+", stripped):
            items.append(clean_text(re.sub(r"^\d+\.\s+", "", stripped)))
        elif re.match(r"^\*\*\d+\.\*\*", stripped):
            # "**1.** Texto del ejercicio."
            cleaned = re.sub(r"^\*\*\d+\.\*\*\s*", "", stripped)
            items.append(clean_text(cleaned))
    return [item for item in items if item]


def parse_markdown_table(markdown_text: str) -> list[list[str]]:
    """Convierte una tabla markdown en filas (sin separador)."""
    rows: list[list[str]] = []
    for line in markdown_text.splitlines():
        stripped = line.strip()
        if not stripped.startswith("|"):
            continue
        row = [clean_text(cell.strip()) for cell in stripped.strip("|").split("|")]
        if all(set(cell) <= set("-: ") for cell in row):
            continue
        rows.append(row)
    return rows


def topics_from_section(body: str) -> list[str]:
    """Extrae los temas como bullets. v3 los presenta como tabla con columna 'Tema'."""
    rows = parse_markdown_table(body)
    if len(rows) >= 2:
        header = [h.lower() for h in rows[0]]
        # Encuentra columna "Tema".
        try:
            tema_idx = next(i for i, h in enumerate(header) if "tema" in h)
        except StopIteration:
            tema_idx = 1 if len(rows[0]) > 1 else 0
        return [row[tema_idx] for row in rows[1:] if len(row) > tema_idx and row[tema_idx]]
    return list_items(body)


def definitions_from_section(body: str) -> list[str]:
    """Extrae nombres + primera línea de definición del bloque de Definiciones."""
    # v3 usa formato:
    #   **Nombre**
    #   : Definición primera línea (puede continuar).
    definitions: list[str] = []
    lines = body.splitlines()
    i = 0
    while i < len(lines):
        line = lines[i].strip()
        m = re.match(r"^\*\*(.+?)\*\*\s*$", line)
        if m and i + 1 < len(lines):
            next_line = lines[i + 1].strip()
            if next_line.startswith(":"):
                desc = clean_text(next_line.lstrip(":").strip())
                # Toma solo la primera oración (hasta el primer punto seguido de espacio).
                first_sentence = re.split(r"(?<=\.)\s", desc, maxsplit=1)[0]
                definitions.append(f"{clean_text(m.group(1))}: {first_sentence}")
                i += 2
                continue
        i += 1
    return definitions


def references_from_section(body: str) -> list[str]:
    """Extrae bullets de referencias."""
    return list_items(body)


def exercises_from_section(body: str) -> list[str]:
    """Extrae enunciados de ejercicios. Toma la primera oración de cada uno."""
    items = list_items(body)
    short: list[str] = []
    for item in items:
        first_sentence = re.split(r"(?<=\.)\s", item, maxsplit=1)[0]
        short.append(first_sentence)
    return short


# ---------------------------------------------------------------------------
# Parser de notebook.ipynb
# ---------------------------------------------------------------------------


def parse_notebook(nb_path: Path) -> dict:
    """Parsea un notebook para extraer la primera celda de código y su intro."""
    data = {
        "n_cells": 0,
        "n_code_cells": 0,
        "first_code": None,
        "first_code_intro": "",
    }
    if not nb_path.exists():
        return data

    cells: list[dict]
    if nbformat is not None:
        try:
            nb = nbformat.read(str(nb_path), as_version=4)
            cells = nb.get("cells", [])
        except Exception:
            cells = []
    else:
        try:
            raw = json.loads(nb_path.read_text(encoding="utf-8"))
            cells = raw.get("cells", [])
        except Exception:
            cells = []

    data["n_cells"] = len(cells)
    data["n_code_cells"] = sum(1 for c in cells if c.get("cell_type") == "code")

    last_markdown_source = ""
    for cell in cells:
        ctype = cell.get("cell_type")
        src = cell.get("source", "")
        if isinstance(src, list):
            src = "".join(src)
        if ctype == "markdown":
            last_markdown_source = src
        elif ctype == "code":
            code_lines = src.splitlines()
            if len(code_lines) > 30:
                code_lines = code_lines[:30] + ["# ... (truncado)"]
            data["first_code"] = "\n".join(code_lines).strip()
            intro_lines = [
                clean_text(line)
                for line in last_markdown_source.splitlines()
                if line.strip() and not line.strip().startswith("#")
            ]
            data["first_code_intro"] = " ".join(intro_lines[:3])[:500]
            break

    return data


# ---------------------------------------------------------------------------
# Construcción del PDF (markdown -> render_markdown_text)
# ---------------------------------------------------------------------------


def available_support_files(class_dir: Path) -> list[str]:
    """Lista archivos complementarios reales presentes en la clase."""
    preferred = ["notebook.ipynb", "homework.ipynb", "soluciones.ipynb", "quiz.json"]
    return [name for name in preferred if (class_dir / name).exists()]


def build_pdf_markdown(class_dir: Path, readme_data: dict, nb_data: dict) -> str:
    """Construye el markdown completo para alimentar render_markdown_text."""
    parts: list[str] = [
        f"# {readme_data['title']}",
        "",
    ]
    if readme_data["subtitle"]:
        parts.extend([f"> {readme_data['subtitle']}", ""])

    for h2, body in readme_data["sections"]:
        body = body.strip()
        if not body:
            continue
        parts.append(f"## {h2}")
        parts.append("")
        parts.append(body)
        parts.append("")

    if nb_data.get("first_code"):
        parts.append("## 💻 Apéndice: notebook (primer bloque)")
        parts.append("")
        intro = nb_data.get("first_code_intro") or (
            "Primera celda ejecutable del notebook de la clase."
        )
        parts.append(intro)
        parts.append("")
        parts.append("```python")
        parts.append(nb_data["first_code"])
        parts.append("```")
        parts.append("")

    support = available_support_files(class_dir)
    if support:
        parts.append("## 📦 Archivos complementarios")
        parts.append("")
        parts.extend(f"- `{name}`" for name in support)
        parts.append("")

    return "\n".join(parts).rstrip() + "\n"


# ---------------------------------------------------------------------------
# Helpers PPTX (paleta, paneles, texto, slides reutilizables)
# ---------------------------------------------------------------------------


def add_background(slide) -> None:
    """Pinta el fondo y la barra superior consistente del deck."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_BG

    top_bar = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, Inches(13.333), Inches(0.45)
    )
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = COLOR_HEADER_BAR
    top_bar.line.color.rgb = COLOR_HEADER_BAR

    accent_bar = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, Inches(0.45), Inches(13.333), Inches(0.08)
    )
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = COLOR_ACCENT
    accent_bar.line.color.rgb = COLOR_ACCENT


def add_panel(slide, x, y, width, height, fill_color=COLOR_PANEL):
    """Dibuja un panel claro con borde suave."""
    panel = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, width, height
    )
    panel.fill.solid()
    panel.fill.fore_color.rgb = fill_color
    panel.line.color.rgb = COLOR_PANEL_BORDER
    panel.line.width = Pt(1.2)
    return panel


def style_paragraph(
    paragraph,
    *,
    font_name: str,
    font_size: int,
    color: RGBColor,
    bold: bool = False,
) -> None:
    """Aplica estilo de párrafo + runs (algunos visores ignoran solo el párrafo)."""
    paragraph.font.name = font_name
    paragraph.font.size = Pt(font_size)
    paragraph.font.color.rgb = color
    paragraph.font.bold = bold
    for run in paragraph.runs:
        run.font.name = font_name
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.bold = bold


def add_title(slide, title: str, subtitle: str = "") -> None:
    """Dibuja el encabezado estándar de cada slide."""
    add_background(slide)

    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.8), Inches(11.9), Inches(0.8))
    frame = title_box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    paragraph = frame.paragraphs[0]
    paragraph.text = title
    style_paragraph(
        paragraph,
        font_name="Segoe UI Semibold",
        font_size=23,
        color=COLOR_TEXT,
        bold=True,
    )

    if subtitle:
        subtitle_box = slide.shapes.add_textbox(
            Inches(0.72), Inches(1.42), Inches(11.9), Inches(0.55)
        )
        sframe = subtitle_box.text_frame
        sframe.clear()
        sframe.word_wrap = True
        sp = sframe.paragraphs[0]
        sp.text = subtitle
        style_paragraph(sp, font_name="Segoe UI", font_size=11, color=COLOR_MUTED)


def add_bullet_slide(slide, title: str, bullets: list[str], subtitle: str = "") -> None:
    """Slide de bullets en un panel grande."""
    add_title(slide, title, subtitle)
    add_panel(slide, Inches(0.8), Inches(1.95), Inches(11.7), Inches(4.9))
    box = slide.shapes.add_textbox(Inches(0.95), Inches(2.05), Inches(11.4), Inches(4.7))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True

    if not bullets:
        bullets = ["Ver README de la clase."]

    for index, bullet in enumerate(bullets):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = f"• {clean_text(bullet)}"
        paragraph.level = 0
        paragraph.space_after = Pt(10)
        style_paragraph(
            paragraph,
            font_name="Segoe UI",
            font_size=18 if len(bullets) <= 4 else 16,
            color=COLOR_TEXT,
        )


def add_two_column_slide(
    slide,
    title: str,
    left_title: str,
    left_values: list[str],
    right_title: str,
    right_values: list[str],
    subtitle: str = "",
) -> None:
    """Dos paneles laterales con header + bullets."""
    add_title(slide, title, subtitle)

    for x, header, values in [
        (Inches(0.8), left_title, left_values or ["Ver README de la clase."]),
        (Inches(6.8), right_title, right_values or ["Ver README de la clase."]),
    ]:
        add_panel(slide, x, Inches(1.95), Inches(5.7), Inches(4.9))

        header_box = slide.shapes.add_textbox(
            x + Inches(0.2), Inches(2.1), Inches(5.3), Inches(0.4)
        )
        header_frame = header_box.text_frame
        header_frame.clear()
        hp = header_frame.paragraphs[0]
        hp.text = header
        style_paragraph(
            hp,
            font_name="Segoe UI Semibold",
            font_size=16,
            color=COLOR_ACCENT,
            bold=True,
        )

        body_box = slide.shapes.add_textbox(
            x + Inches(0.2), Inches(2.55), Inches(5.3), Inches(4.0)
        )
        body_frame = body_box.text_frame
        body_frame.clear()
        body_frame.word_wrap = True
        for index, value in enumerate(values):
            paragraph = body_frame.paragraphs[0] if index == 0 else body_frame.add_paragraph()
            paragraph.text = f"• {clean_text(value)}"
            style_paragraph(
                paragraph,
                font_name="Segoe UI",
                font_size=14 if len(values) > 4 else 15,
                color=COLOR_TEXT,
            )
            paragraph.space_after = Pt(8)


def add_code_slide(slide, title: str, intro: str, code: str) -> None:
    """Slide con un bloque de código monospace en panel oscuro."""
    add_title(slide, title, intro)
    code_panel = add_panel(
        slide,
        Inches(0.8),
        Inches(2.0),
        Inches(11.7),
        Inches(4.9),
        fill_color=COLOR_CODE,
    )
    code_panel.fill.solid()
    code_panel.fill.fore_color.rgb = COLOR_CODE
    code_panel.line.color.rgb = COLOR_HEADER_BAR

    code_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.2), Inches(4.5))
    frame = code_box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.TOP

    # Limitar a 20 líneas máx.
    code_lines = code.splitlines()
    if len(code_lines) > 20:
        code_lines = code_lines[:20] + ["# ... (truncado)"]

    for idx, line in enumerate(code_lines):
        paragraph = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        # No usar clean_text — el código debe conservar `_`, `*`, etc.
        paragraph.text = line if line else " "
        style_paragraph(
            paragraph,
            font_name="Courier New",
            font_size=12,
            color=COLOR_CODE_TEXT,
        )


# ---------------------------------------------------------------------------
# Construcción del deck PPTX
# ---------------------------------------------------------------------------


def create_presentation(readme_data: dict, nb_data: dict, output_path: Path) -> None:
    """Construye el deck `.pptx` con ~7-8 slides desde la pauta v3."""
    title = readme_data["title"]
    subtitle = readme_data["subtitle"]

    objective = get_section(readme_data, "Objetivo")
    outcomes_body = get_section(readme_data, "Resultados de aprendizaje", "Resultados")
    outcomes = list_items(outcomes_body)
    topics_body = get_section(readme_data, "Temas")
    topics = topics_from_section(topics_body)
    definitions_body = get_section(readme_data, "Definiciones")
    definitions = definitions_from_section(definitions_body)
    exercises_body = get_section(readme_data, "Ejercicios")
    exercises = exercises_from_section(exercises_body)
    homework_body = get_section(readme_data, "Homework")
    homework = list_items(homework_body)
    if not homework and homework_body:
        # El homework v3 suele ser un párrafo descriptivo + criterio.
        first_para = next(
            (p.strip() for p in homework_body.split("\n\n") if p.strip()), ""
        )
        if first_para:
            homework = [clean_text(first_para)]
    references_body = get_section(readme_data, "Referencias")
    references = references_from_section(references_body)
    next_body = get_section(readme_data, "Siguiente clase")

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # Slide 1 — Portada.
    slide = prs.slides.add_slide(blank)
    add_title(slide, title, subtitle)
    add_panel(slide, Inches(0.8), Inches(1.95), Inches(11.7), Inches(4.9))
    intro_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.15), Inches(11.3), Inches(4.5))
    frame = intro_box.text_frame
    frame.clear()
    frame.word_wrap = True
    paragraph = frame.paragraphs[0]
    paragraph.text = objective.strip().split("\n\n")[0] if objective else title
    style_paragraph(
        paragraph,
        font_name="Segoe UI Semibold",
        font_size=22,
        color=COLOR_TEXT,
        bold=True,
    )
    paragraph.space_after = Pt(18)
    if outcomes:
        sub = frame.add_paragraph()
        sub.text = "Al cierre: " + outcomes[0]
        style_paragraph(sub, font_name="Segoe UI", font_size=15, color=COLOR_MUTED)

    # Slide 2 — Objetivo + Resultados.
    slide = prs.slides.add_slide(blank)
    add_two_column_slide(
        slide,
        "🎯 Objetivo · 📚 Resultados",
        "Objetivo",
        [objective.strip().split("\n\n")[0]] if objective else [],
        "Resultados de aprendizaje",
        outcomes[:5],
        subtitle,
    )

    # Slide 3 — Temas.
    slide = prs.slides.add_slide(blank)
    add_bullet_slide(slide, "🗺️ Temas", topics[:7], "Recorrido de la sesión.")

    # Slide 4 — Definiciones clave.
    slide = prs.slides.add_slide(blank)
    add_bullet_slide(
        slide,
        "📖 Definiciones clave",
        definitions[:6],
        "Conceptos a fijar antes de la práctica.",
    )

    # Slide 5 — Código de muestra.
    if nb_data.get("first_code"):
        slide = prs.slides.add_slide(blank)
        intro = nb_data.get("first_code_intro") or "Primer bloque ejecutable del notebook."
        add_code_slide(slide, "💻 Código de muestra", intro, nb_data["first_code"])

    # Slide 6 — Ejercicios + Homework.
    slide = prs.slides.add_slide(blank)
    add_two_column_slide(
        slide,
        "🧪 Ejercicios · 📝 Homework",
        "Ejercicios",
        exercises[:5],
        "Homework verificable",
        homework[:5],
        "Práctica guiada + entrega.",
    )

    # Slide 7 — Referencias.
    slide = prs.slides.add_slide(blank)
    add_bullet_slide(
        slide,
        "🔗 Referencias",
        references[:6],
        "Bibliografía y docs oficiales.",
    )

    # Slide 8 — Cierre.
    slide = prs.slides.add_slide(blank)
    closing_text = "Pasar a la práctica con el notebook."
    if next_body:
        first_line = next(
            (line.strip() for line in next_body.splitlines() if line.strip()),
            "",
        )
        if first_line:
            closing_text = clean_text(first_line)
    add_title(slide, "Siguiente paso →", "Cierre de la clase.")
    add_panel(slide, Inches(0.8), Inches(1.95), Inches(11.7), Inches(4.9))
    box = slide.shapes.add_textbox(Inches(1.0), Inches(3.0), Inches(11.3), Inches(2.5))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    p = frame.paragraphs[0]
    p.text = closing_text
    style_paragraph(
        p, font_name="Segoe UI Semibold", font_size=24, color=COLOR_TEXT, bold=True
    )

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)


# ---------------------------------------------------------------------------
# Orquestación
# ---------------------------------------------------------------------------


def class_output_name(class_dir: Path, suffix: str) -> str:
    """Compone `clase-<basename>-<suffix>` (solo basename, sin slug de parte)."""
    return f"clase-{class_dir.name}-{suffix}"


def generate_assets(class_dirs: list[Path], force: bool = False) -> tuple[int, int]:
    """Genera PDF + PPTX para cada clase. Devuelve (ok_count, error_count)."""
    ok = 0
    errors = 0

    for class_dir in class_dirs:
        readme_path = class_dir / "README.md"
        nb_path = class_dir / "notebook.ipynb"

        pdf_filename = class_output_name(class_dir, "guia-explicativa.pdf")
        pptx_filename = class_output_name(class_dir, "presentacion.pptx")

        pdf_mirror = PDF_OUTPUT_DIR / pdf_filename
        pptx_mirror = PPTX_OUTPUT_DIR / pptx_filename
        pdf_local = class_dir / pdf_filename
        pptx_local = class_dir / pptx_filename

        if (
            not force
            and pdf_mirror.exists()
            and pptx_mirror.exists()
            and pdf_local.exists()
            and pptx_local.exists()
        ):
            print(f"[SKIP] {class_dir.relative_to(BASE_DIR)} (ya existen, usar --force)")
            ok += 1
            continue

        try:
            readme_data = parse_class_readme(readme_path)
            nb_data = parse_notebook(nb_path)

            # PDF.
            markdown_text = build_pdf_markdown(class_dir, readme_data, nb_data)
            pdf_local.parent.mkdir(parents=True, exist_ok=True)
            render_markdown_text(
                markdown_text=markdown_text,
                output_path=pdf_local,
                title=readme_data["title"],
                subtitle=readme_data["subtitle"]
                or "Guía construida desde el contenido real de la clase.",
                style="print",
            )

            # PPTX.
            create_presentation(readme_data, nb_data, pptx_local)

            # Mirror a docs/.
            pdf_mirror.parent.mkdir(parents=True, exist_ok=True)
            pptx_mirror.parent.mkdir(parents=True, exist_ok=True)
            pdf_mirror.write_bytes(pdf_local.read_bytes())
            pptx_mirror.write_bytes(pptx_local.read_bytes())

            print(f"[OK] {pdf_local.relative_to(BASE_DIR)}")
            print(f"[OK] {pptx_local.relative_to(BASE_DIR)}")
            print(f"[OK] {pdf_mirror.relative_to(BASE_DIR)}")
            print(f"[OK] {pptx_mirror.relative_to(BASE_DIR)}")
            ok += 1
        except Exception as exc:
            errors += 1
            print(
                f"[ERROR] {class_dir.relative_to(BASE_DIR)}: {exc.__class__.__name__}: {exc}",
                file=sys.stderr,
            )
            traceback.print_exc(file=sys.stderr)

    return ok, errors


def main() -> None:
    """CLI principal — genera activos PDF + PPTX para clases v3."""
    parser = argparse.ArgumentParser(
        description="Genera PDFs y presentaciones por clase para el currículo v3."
    )
    parser.add_argument(
        "--parte",
        type=int,
        default=None,
        help="Filtra a una sola parte (0-8).",
    )
    parser.add_argument(
        "--clase",
        type=str,
        default=None,
        help=(
            "Filtra a una sola clase. Acepta basename "
            "(ej. '006-python-...') o slug 'parte-0-.../006-...'."
        ),
    )
    parser.add_argument(
        "--force",
        action="store_true",
        help="Regenera aunque los archivos ya existan.",
    )
    args = parser.parse_args()

    all_dirs = discover_class_dirs()
    dirs = filter_dirs(all_dirs, parte=args.parte, clase=args.clase)

    if not dirs:
        raise SystemExit(
            "No se encontró ninguna clase con los filtros indicados. "
            f"(Total descubierto: {len(all_dirs)})"
        )

    print(f"Procesando {len(dirs)} clase(s) de {len(all_dirs)} totales...")
    ok, errors = generate_assets(dirs, force=args.force)
    print(f"\nResumen: {ok} OK · {errors} errores · {len(dirs)} totales.")
    if errors:
        sys.exit(1)


if __name__ == "__main__":
    main()
