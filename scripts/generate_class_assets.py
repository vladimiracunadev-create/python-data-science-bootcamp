"""Genera PDFs y presentaciones por clase usando el contenido real del repo.

Responsabilidades:
    1. Leer `README.md`, `slides.md`, `teoria.md`, `ejercicios.md` y
       `homework.md` dentro de cada carpeta de clase.
    2. Construir una guía PDF que reúna la información completa del módulo.
    3. Crear un deck `.pptx` que traduzca ese contenido a una secuencia visual.

Qué resuelve:
    Evita presentar clases con una plantilla genérica. Cada salida nace del
    material específico que ya existe en la clase correspondiente.
"""

from __future__ import annotations

import argparse
import re
import unicodedata
from pathlib import Path

from generar_pdf_documento import render_markdown_text
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR
from pptx.util import Inches, Pt

BASE_DIR = Path(__file__).resolve().parents[1]
CLASSES_DIR = BASE_DIR / "classes"
PDF_OUTPUT_DIR = BASE_DIR / "docs" / "pdfs" / "classes"
PPTX_OUTPUT_DIR = BASE_DIR / "docs" / "presentaciones" / "classes"

SOURCE_FILES = [
    ("README.md", "📘 Ficha de clase"),
    ("slides.md", "🖥️ Guion de diapositivas"),
    ("teoria.md", "🧠 Desarrollo teórico"),
    ("ejercicios.md", "🧪 Ejercicios y práctica"),
    ("homework.md", "📝 Tarea y cierre"),
]

COLOR_BG = RGBColor(15, 23, 42)
COLOR_PANEL = RGBColor(17, 24, 39)
COLOR_ACCENT = RGBColor(34, 197, 94)
COLOR_TEXT = RGBColor(241, 245, 249)
COLOR_MUTED = RGBColor(148, 163, 184)
COLOR_CODE = RGBColor(2, 6, 23)

EMOJI_RE = re.compile(r"[\U0001F300-\U0001FAFF\u2600-\u27BF\u2300-\u23FF]")


def read_text(path: Path) -> str:
    """Lee un archivo markdown con UTF-8."""
    return path.read_text(encoding="utf-8")


def clean_text(text: str) -> str:
    """Quita markdown inline simple y emojis para reutilizar texto en PPTX."""
    clean = text
    clean = re.sub(r"\[(.+?)\]\(.+?\)", r"\1", clean)
    clean = re.sub(r"`(.+?)`", r"\1", clean)
    clean = re.sub(r"\*\*(.+?)\*\*", r"\1", clean)
    clean = re.sub(r"\*(.+?)\*", r"\1", clean)
    clean = EMOJI_RE.sub("", clean)
    clean = clean.replace("\u200d", "")
    clean = clean.replace("\ufe0f", "")
    return clean.strip()


def normalize_heading(text: str) -> str:
    """Normaliza un heading para buscarlo por contenido semántico."""
    plain = clean_text(text).lower()
    plain = unicodedata.normalize("NFD", plain)
    plain = "".join(char for char in plain if unicodedata.category(char) != "Mn")
    return re.sub(r"\s+", " ", plain).strip()


def extract_h1(markdown_text: str, fallback: str) -> str:
    """Obtiene el H1 principal del documento."""
    for line in markdown_text.splitlines():
        if line.startswith("# "):
            return clean_text(line[2:])
    return fallback


def strip_first_h1(markdown_text: str) -> str:
    """Quita el primer H1 para evitar duplicarlo en PDF/PPTX."""
    lines = markdown_text.splitlines()
    for index, line in enumerate(lines):
        if line.startswith("# "):
            return "\n".join(lines[index + 1 :]).lstrip()
    return markdown_text


def split_h2_sections(markdown_text: str) -> list[tuple[str, str]]:
    """Divide un markdown en secciones a nivel `##` preservando orden."""
    body = strip_first_h1(markdown_text)
    sections: list[tuple[str, list[str]]] = []
    current_title: str | None = None
    current_lines: list[str] = []

    for line in body.splitlines():
        if line.startswith("## "):
            if current_title is not None:
                sections.append((current_title, current_lines[:]))
            current_title = line[3:].strip()
            current_lines = []
        else:
            current_lines.append(line)

    if current_title is not None:
        sections.append((current_title, current_lines[:]))

    return [(title, "\n".join(lines).strip()) for title, lines in sections]


def find_section(markdown_text: str, needle: str) -> str:
    """Busca una sección H2 por coincidencia parcial normalizada."""
    target = normalize_heading(needle)
    for title, body in split_h2_sections(markdown_text):
        if target in normalize_heading(title):
            return body
    return ""


def list_items(markdown_text: str) -> list[str]:
    """Extrae bullets y listas numeradas de una sección markdown."""
    items: list[str] = []
    for line in markdown_text.splitlines():
        stripped = line.strip()
        if stripped.startswith("- ") or stripped.startswith("* "):
            items.append(clean_text(stripped[2:]))
        elif re.match(r"^\d+\. ", stripped):
            items.append(clean_text(re.sub(r"^\d+\. ", "", stripped)))
    return [item for item in items if item]


def first_paragraphs(markdown_text: str, limit: int = 2) -> list[str]:
    """Extrae hasta `limit` párrafos limpios para síntesis visual."""
    paragraphs: list[str] = []
    current: list[str] = []
    in_code = False

    for line in markdown_text.splitlines():
        stripped = line.strip()
        if stripped.startswith("```"):
            in_code = not in_code
            continue
        if in_code or stripped.startswith("|"):
            continue
        if stripped.startswith("### "):
            if current:
                paragraphs.append(clean_text(" ".join(current)))
                current = []
            if len(paragraphs) >= limit:
                break
            paragraphs.append(clean_text(stripped[4:]))
            if len(paragraphs) >= limit:
                break
            continue
        if not stripped:
            if current:
                paragraphs.append(clean_text(" ".join(current)))
                current = []
            if len(paragraphs) >= limit:
                break
            continue
        if stripped.startswith("- ") or stripped.startswith("* ") or re.match(r"^\d+\. ", stripped):
            continue
        current.append(stripped)

    if current and len(paragraphs) < limit:
        paragraphs.append(clean_text(" ".join(current)))

    return [paragraph for paragraph in paragraphs if paragraph][:limit]


def parse_markdown_table(markdown_text: str) -> list[list[str]]:
    """Convierte una tabla markdown en filas simples."""
    rows: list[list[str]] = []
    for line in markdown_text.splitlines():
        stripped = line.strip()
        if not stripped.startswith("|"):
            continue
        row = [clean_text(cell.strip()) for cell in stripped.strip("|").split("|")]
        if all(set(cell) <= set("-: ") for cell in row):
            continue
        rows.append(row)
    return rows


def table_rows_as_bullets(markdown_text: str) -> list[str]:
    """Transforma la tabla de ruta de sesión en bullets legibles."""
    rows = parse_markdown_table(markdown_text)
    if len(rows) < 2:
        return []

    bullets: list[str] = []
    for row in rows[1:]:
        if len(row) >= 4:
            bullets.append(f"{row[0]}: {row[1]} | {row[2]} | Evidencia: {row[3]}")
    return bullets


def first_code_block(markdown_text: str) -> tuple[str, str, str]:
    """Extrae el primer bloque de código y su contexto cercano."""
    current_h3 = "Bloque principal"
    intro_lines: list[str] = []
    code_lines: list[str] = []
    capture_intro = False
    in_code = False

    for line in markdown_text.splitlines():
        stripped = line.strip()
        if stripped.startswith("### "):
            current_h3 = clean_text(stripped[4:])
            intro_lines = []
            capture_intro = True
            continue
        if stripped.startswith("```"):
            if not in_code:
                in_code = True
                continue
            return current_h3, clean_text(" ".join(intro_lines)), "\n".join(code_lines).strip()
        if in_code:
            code_lines.append(line.rstrip())
        elif capture_intro and stripped:
            intro_lines.append(stripped)

    return current_h3, clean_text(" ".join(intro_lines)), "\n".join(code_lines).strip()


def available_support_files(class_dir: Path) -> list[str]:
    """Lista archivos relevantes del módulo que no entran en la cadena de markdown."""
    preferred = ["quiz.json", "notebook.ipynb", "soluciones.ipynb"]
    files = [name for name in preferred if (class_dir / name).exists()]
    return files


def build_class_markdown(class_dir: Path) -> str:
    """Construye la guía PDF a partir del contenido real de la clase."""
    readme_text = read_text(class_dir / "README.md")
    title = extract_h1(readme_text, class_dir.name)

    sections: list[str] = [
        f"# {title}",
        "",
        "> 📘 Guía integral derivada del contenido real del módulo.",
        "",
    ]

    for filename, label in SOURCE_FILES:
        path = class_dir / filename
        if not path.exists():
            continue
        body = strip_first_h1(read_text(path)).strip()
        if not body:
            continue
        sections.extend([f"## {label}", "", body, ""])

    support = available_support_files(class_dir)
    if support:
        sections.extend(
            [
                "## 📦 Archivos complementarios",
                "",
                *[f"- `{name}`" for name in support],
                "",
            ]
        )

    return "\n".join(sections).rstrip() + "\n"


def class_context(class_dir: Path) -> dict[str, str | list[str]]:
    """Prepara el contexto reutilizable para construir el deck PPTX."""
    readme_text = read_text(class_dir / "README.md")
    slides_text = read_text(class_dir / "slides.md")
    theory_text = read_text(class_dir / "teoria.md")
    exercises_text = read_text(class_dir / "ejercicios.md")
    homework_text = read_text(class_dir / "homework.md")

    title = extract_h1(readme_text, class_dir.name)
    objective = " ".join(first_paragraphs(find_section(readme_text, "Objetivo"), limit=1))
    duration = " ".join(list_items(find_section(readme_text, "Duración sugerida"))[:1])
    dataset = " ".join(list_items(find_section(readme_text, "Dataset base"))[:1])
    outcomes = list_items(find_section(readme_text, "Resultados esperados"))
    topics = list_items(find_section(readme_text, "Temas clave"))
    idea = " ".join(first_paragraphs(find_section(readme_text, "Idea central"), limit=1))
    note = " ".join(first_paragraphs(find_section(readme_text, "Nota para el docente"), limit=1))
    route = table_rows_as_bullets(find_section(slides_text, "Ruta de la sesión"))
    key_points = list_items(find_section(slides_text, "Puntos que deben quedar claros"))
    closure = " ".join(first_paragraphs(find_section(slides_text, "Cierre esperado"), limit=1))
    why = " ".join(first_paragraphs(find_section(theory_text, "Por qué importa"), limit=1))
    theory_idea = " ".join(first_paragraphs(find_section(theory_text, "Idea central"), limit=1))
    errors = list_items(find_section(theory_text, "Errores frecuentes"))
    connection = " ".join(first_paragraphs(find_section(theory_text, "Conexión"), limit=1))
    code_title, code_intro, code = first_code_block(theory_text)
    guided_work = list_items(find_section(exercises_text, "Trabajo guiado"))
    self_checks = list_items(find_section(exercises_text, "Criterios de autocorrección"))
    homework_assignment = first_paragraphs(find_section(homework_text, "Encargo"), limit=1)
    deliverables = list_items(find_section(homework_text, "Entregables"))
    final_checks = list_items(find_section(homework_text, "Autoevaluación final"))
    support = available_support_files(class_dir)

    return {
        "title": title,
        "objective": objective,
        "duration": duration,
        "dataset": dataset,
        "outcomes": outcomes,
        "topics": topics,
        "idea": idea,
        "note": note,
        "route": route,
        "key_points": key_points,
        "closure": closure,
        "why": why,
        "theory_idea": theory_idea,
        "errors": errors,
        "connection": connection,
        "code_title": code_title,
        "code_intro": code_intro,
        "code": code,
        "guided_work": guided_work,
        "self_checks": self_checks,
        "homework_assignment": homework_assignment,
        "deliverables": deliverables,
        "final_checks": final_checks,
        "support": support,
    }


def add_background(slide) -> None:
    """Pinta el fondo general del slide."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_BG

    top_bar = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        0,
        0,
        Inches(13.333),
        Inches(0.45),
    )
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = COLOR_PANEL
    top_bar.line.color.rgb = COLOR_PANEL

    accent_bar = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        0,
        Inches(0.45),
        Inches(13.333),
        Inches(0.08),
    )
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = COLOR_ACCENT
    accent_bar.line.color.rgb = COLOR_ACCENT


def add_title(slide, title: str, subtitle: str = "") -> None:
    """Dibuja encabezado consistente para todos los slides."""
    add_background(slide)

    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.8), Inches(11.9), Inches(0.8))
    frame = title_box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    paragraph = frame.paragraphs[0]
    paragraph.text = title
    paragraph.font.name = "Segoe UI Semibold"
    paragraph.font.size = Pt(23)
    paragraph.font.color.rgb = COLOR_TEXT

    if subtitle:
        subtitle_box = slide.shapes.add_textbox(Inches(0.72), Inches(1.42), Inches(11.9), Inches(0.55))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.clear()
        subtitle_frame.word_wrap = True
        paragraph = subtitle_frame.paragraphs[0]
        paragraph.text = subtitle
        paragraph.font.name = "Segoe UI"
        paragraph.font.size = Pt(11)
        paragraph.font.color.rgb = COLOR_MUTED


def add_bullet_slide(slide, title: str, bullets: list[str], subtitle: str = "") -> None:
    """Agrega una slide de bullets basada en contenido del repo."""
    add_title(slide, title, subtitle)
    box = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(11.7), Inches(4.7))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True

    for index, bullet in enumerate(bullets):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = clean_text(bullet)
        paragraph.level = 0
        paragraph.space_after = Pt(10)
        paragraph.font.name = "Segoe UI"
        paragraph.font.size = Pt(20 if len(bullets) <= 4 else 18)
        paragraph.font.color.rgb = COLOR_TEXT


def add_two_column_slide(slide, title: str, left_title: str, left_values: list[str], right_title: str, right_values: list[str], subtitle: str = "") -> None:
    """Dibuja dos paneles laterales para resumir objetivos, temas o entregables."""
    add_title(slide, title, subtitle)

    for x, header, values in [
        (Inches(0.8), left_title, left_values),
        (Inches(6.8), right_title, right_values),
    ]:
        panel = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            x,
            Inches(1.95),
            Inches(5.1),
            Inches(4.5),
        )
        panel.fill.solid()
        panel.fill.fore_color.rgb = COLOR_PANEL
        panel.line.color.rgb = COLOR_PANEL

        header_box = slide.shapes.add_textbox(x + Inches(0.2), Inches(2.15), Inches(4.7), Inches(0.4))
        header_frame = header_box.text_frame
        header_frame.clear()
        paragraph = header_frame.paragraphs[0]
        paragraph.text = header
        paragraph.font.name = "Segoe UI Semibold"
        paragraph.font.size = Pt(16)
        paragraph.font.color.rgb = COLOR_ACCENT

        body_box = slide.shapes.add_textbox(x + Inches(0.2), Inches(2.6), Inches(4.7), Inches(3.5))
        body_frame = body_box.text_frame
        body_frame.clear()
        body_frame.word_wrap = True
        for index, value in enumerate(values):
            paragraph = body_frame.paragraphs[0] if index == 0 else body_frame.add_paragraph()
            paragraph.text = clean_text(value)
            paragraph.font.name = "Segoe UI"
            paragraph.font.size = Pt(17)
            paragraph.font.color.rgb = COLOR_TEXT
            paragraph.space_after = Pt(8)


def add_code_slide(slide, title: str, intro: str, code: str) -> None:
    """Agrega un slide de código basado en el bloque real de teoría."""
    add_title(slide, title, intro)

    code_panel = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.8),
        Inches(2.0),
        Inches(11.7),
        Inches(4.7),
    )
    code_panel.fill.solid()
    code_panel.fill.fore_color.rgb = COLOR_CODE
    code_panel.line.color.rgb = COLOR_PANEL

    code_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.22), Inches(11.2), Inches(4.2))
    frame = code_box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.TOP
    paragraph = frame.paragraphs[0]
    paragraph.text = code
    paragraph.font.name = "Courier New"
    paragraph.font.size = Pt(12)
    paragraph.font.color.rgb = COLOR_TEXT


def create_presentation(class_dir: Path, output_path: Path) -> None:
    """Construye el deck `.pptx` reutilizando el contenido específico de la clase."""
    ctx = class_context(class_dir)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    slide = prs.slides.add_slide(blank)
    subtitle_parts = [ctx["duration"], ctx["dataset"]]
    subtitle = " | ".join(part for part in subtitle_parts if part)
    add_title(slide, str(ctx["title"]), subtitle)
    intro_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.12), Inches(11.7), Inches(3.2))
    intro_frame = intro_box.text_frame
    intro_frame.clear()
    intro_frame.word_wrap = True
    paragraph = intro_frame.paragraphs[0]
    paragraph.text = str(ctx["objective"])
    paragraph.font.name = "Segoe UI Semibold"
    paragraph.font.size = Pt(27)
    paragraph.font.color.rgb = COLOR_TEXT
    paragraph.space_after = Pt(18)
    if ctx["idea"]:
        paragraph = intro_frame.add_paragraph()
        paragraph.text = str(ctx["idea"])
        paragraph.font.name = "Segoe UI"
        paragraph.font.size = Pt(18)
        paragraph.font.color.rgb = COLOR_MUTED

    slide = prs.slides.add_slide(blank)
    add_two_column_slide(
        slide,
        "🎯 Foco del módulo",
        "Resultados esperados",
        list(ctx["outcomes"]) or ["Revisar la ficha del módulo."],
        "Temas clave",
        list(ctx["topics"]) or ["Revisar los temas clave en README.md."],
        "Resumen tomado desde README.md.",
    )

    if ctx["route"]:
        slide = prs.slides.add_slide(blank)
        add_bullet_slide(
            slide,
            "🖥️ Ruta sugerida de la sesión",
            list(ctx["route"]),
            "Resumen tomado desde slides.md.",
        )

    key_points = list(ctx["key_points"]) or list(ctx["topics"])
    if key_points or ctx["closure"]:
        bullets = key_points[:]
        if ctx["closure"]:
            bullets.append(f"Cierre esperado: {ctx['closure']}")
        slide = prs.slides.add_slide(blank)
        add_bullet_slide(
            slide,
            "📌 Puntos que deben quedar claros",
            bullets,
            "Contenido derivado del guion de diapositivas.",
        )

    theory_bullets = [value for value in [ctx["theory_idea"], ctx["why"], ctx["connection"]] if value]
    if theory_bullets:
        slide = prs.slides.add_slide(blank)
        add_bullet_slide(
            slide,
            "🧠 Base conceptual del módulo",
            theory_bullets,
            "Síntesis tomada desde teoria.md.",
        )

    if ctx["code"]:
        slide = prs.slides.add_slide(blank)
        add_code_slide(
            slide,
            f"💻 {ctx['code_title']}",
            str(ctx["code_intro"]),
            str(ctx["code"]),
        )
    elif ctx["errors"]:
        slide = prs.slides.add_slide(blank)
        add_bullet_slide(slide, "⚠️ Errores frecuentes", list(ctx["errors"]))

    exercise_left = list(ctx["guided_work"]) or first_paragraphs(find_section(read_text(class_dir / "ejercicios.md"), "Trabajo guiado"), limit=3)
    exercise_right = list(ctx["self_checks"]) or ["Revisar la autocorrección definida en ejercicios.md."]
    slide = prs.slides.add_slide(blank)
    add_two_column_slide(
        slide,
        "🧪 Práctica guiada",
        "Trabajo guiado",
        exercise_left[:4] if exercise_left else ["Explicar la consigna antes de resolver."],
        "Autocorrección",
        exercise_right[:4] if exercise_right else ["Verificar si la salida responde a la pregunta."],
        "Derivado del archivo ejercicios.md.",
    )

    homework_left = list(ctx["deliverables"]) or [*ctx["homework_assignment"]] or ["Revisar el encargo del módulo."]
    homework_right = list(ctx["final_checks"]) or list(ctx["support"]) or ["Revisar homework.md para el cierre."]
    slide = prs.slides.add_slide(blank)
    add_two_column_slide(
        slide,
        "📝 Tarea y evidencia final",
        "Entregables",
        homework_left[:4],
        "Autoevaluación o apoyo",
        homework_right[:4],
        "Construido desde homework.md y archivos complementarios.",
    )

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)


def class_output_name(class_dir: Path, suffix: str) -> str:
    """Compone un nombre de archivo estable usando el slug de la clase."""
    return f"clase-{class_dir.name}-{suffix}"


def class_directories(selected_class: str | None) -> list[Path]:
    """Lista las carpetas de clase ordenadas o una sola si se pide por CLI."""
    directories = sorted(path for path in CLASSES_DIR.iterdir() if path.is_dir())
    if selected_class is None:
        return directories
    return [path for path in directories if path.name.startswith(selected_class)]


def generate_assets(directories: list[Path]) -> None:
    """Genera PDF y PPTX para las clases indicadas."""
    for class_dir in directories:
        pdf_path = PDF_OUTPUT_DIR / f"{class_output_name(class_dir, 'guia-explicativa.pdf')}"
        pptx_path = PPTX_OUTPUT_DIR / f"{class_output_name(class_dir, 'presentacion.pptx')}"

        markdown = build_class_markdown(class_dir)
        title = extract_h1(read_text(class_dir / "README.md"), class_dir.name)

        render_markdown_text(
            markdown_text=markdown,
            output_path=pdf_path,
            title=title,
            subtitle="Guía construida a partir del contenido real del módulo.",
        )
        create_presentation(class_dir, pptx_path)

        print(f"[OK] {pdf_path.relative_to(BASE_DIR)}")
        print(f"[OK] {pptx_path.relative_to(BASE_DIR)}")


def main() -> None:
    """CLI principal para generar activos por clase."""
    parser = argparse.ArgumentParser(description="Genera PDFs y presentaciones por clase desde el contenido real del repo.")
    parser.add_argument("--clase", help="Prefijo de clase, por ejemplo 01 o 12.")
    args = parser.parse_args()

    directories = class_directories(args.clase)
    if not directories:
        raise SystemExit("No se encontró la clase solicitada.")

    generate_assets(directories)


if __name__ == "__main__":
    main()
