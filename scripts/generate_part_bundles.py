"""Genera PDFs y PPTX consolidados por parte y un bundle unificado del curso.

Qué resuelve:
    El currículo v3 tiene 232 clases en 9 partes. Los assets generados por
    `generate_class_assets_v3.py` son uno-por-clase. Este script consolida:

        1. `docs/pdfs/parts/parte-N-slug-completa.pdf` — todas las clases de
           la parte concatenadas en un solo PDF (estilo `print`).
        2. `docs/presentaciones/parts/parte-N-slug-completa.pptx` — todos los
           slides de las clases de la parte en un solo deck.
        3. `docs/pdfs/curso-completo.pdf` — el currículo entero.
        4. `docs/presentaciones/curso-completo.pptx` — slides del curso entero.

    Reutiliza el renderer Markdown→PDF de `scripts/generar_pdf_documento.py`
    y los helpers visuales de PPTX de `generate_class_assets_v3.py`. NO
    modifica ninguno de esos módulos.

Uso:
    python scripts/generate_part_bundles.py            # genera todo (9 + unificado)
    python scripts/generate_part_bundles.py --all
    python scripts/generate_part_bundles.py --parte 0
    python scripts/generate_part_bundles.py --unified
    python scripts/generate_part_bundles.py --skip-pdf
    python scripts/generate_part_bundles.py --skip-pptx
"""

from __future__ import annotations

import argparse
import re
import sys
import traceback
from pathlib import Path

# Reutilizar el renderer PDF y los helpers PPTX existentes del repo.
sys.path.insert(0, str(Path(__file__).resolve().parent))
from generar_pdf_documento import render_markdown_text  # noqa: E402
from generate_class_assets_v3 import (  # noqa: E402
    COLOR_ACCENT,
    COLOR_BG,
    COLOR_CODE,
    COLOR_CODE_TEXT,
    COLOR_HEADER_BAR,
    COLOR_MUTED,
    COLOR_PANEL,
    COLOR_TEXT,
    add_background,
    add_bullet_slide,
    add_code_slide,
    add_panel,
    add_title,
    add_two_column_slide,
    clean_text,
    discover_class_dirs,
    exercises_from_section,
    get_section,
    list_items,
    parse_class_readme,
    parse_notebook,
    style_paragraph,
    topics_from_section,
)
from pptx import Presentation  # noqa: E402
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE  # noqa: E402
from pptx.enum.text import MSO_ANCHOR  # noqa: E402
from pptx.util import Inches, Pt  # noqa: E402

BASE_DIR = Path(__file__).resolve().parents[1]
CLASSES_DIR = BASE_DIR / "classes"
PDF_PARTS_DIR = BASE_DIR / "docs" / "pdfs" / "parts"
PPTX_PARTS_DIR = BASE_DIR / "docs" / "presentaciones" / "parts"
PDF_UNIFIED = BASE_DIR / "docs" / "pdfs" / "curso-completo.pdf"
PPTX_UNIFIED = BASE_DIR / "docs" / "presentaciones" / "curso-completo.pptx"

PARTE_DIR_RE = re.compile(r"^parte-(\d+)-(.+)$")
CLASS_DIR_RE = re.compile(r"^(\d{3})-(.+)$")
NEXT_CLASS_HEADING_RE = re.compile(
    r"^##\s+.*siguiente\s+clase.*$", re.IGNORECASE | re.MULTILINE
)


# ---------------------------------------------------------------------------
# Descubrimiento por parte
# ---------------------------------------------------------------------------


def list_parts() -> list[tuple[int, str, Path]]:
    """Devuelve [(parte_num, slug, dir)] ordenado por número de parte."""
    parts: list[tuple[int, str, Path]] = []
    for parte_dir in sorted(CLASSES_DIR.iterdir()):
        if not parte_dir.is_dir():
            continue
        m = PARTE_DIR_RE.match(parte_dir.name)
        if not m:
            continue
        parts.append((int(m.group(1)), m.group(2), parte_dir))
    parts.sort(key=lambda item: item[0])
    return parts


def classes_in_part(part_dir: Path) -> list[Path]:
    """Devuelve carpetas de clase dentro de una parte, ordenadas por NNN."""
    items: list[tuple[int, Path]] = []
    for class_dir in sorted(part_dir.iterdir()):
        if not class_dir.is_dir():
            continue
        m = CLASS_DIR_RE.match(class_dir.name)
        if not m:
            continue
        if not (class_dir / "README.md").exists():
            continue
        items.append((int(m.group(1)), class_dir))
    items.sort(key=lambda x: x[0])
    return [d for _, d in items]


def read_part_meta(part_dir: Path) -> tuple[str, str]:
    """Lee el README de la parte (si existe) para extraer título y subtítulo."""
    readme = part_dir / "README.md"
    if not readme.exists():
        # Fallback: título humano a partir del slug.
        slug = part_dir.name
        title = slug.replace("-", " ").title()
        return title, ""
    text = readme.read_text(encoding="utf-8")
    title = part_dir.name
    subtitle_parts: list[str] = []
    body_start = 0
    lines = text.splitlines()
    for idx, line in enumerate(lines):
        if line.startswith("# "):
            title = clean_text(line[2:])
            body_start = idx + 1
            break
    cursor = body_start
    while cursor < len(lines):
        stripped = lines[cursor].strip()
        if stripped.startswith(">"):
            subtitle_parts.append(clean_text(stripped.lstrip(">").strip()))
            cursor += 1
        elif stripped == "":
            cursor += 1
        else:
            break
    return title, " · ".join(p for p in subtitle_parts if p)


# ---------------------------------------------------------------------------
# Helpers markdown
# ---------------------------------------------------------------------------


def demote_headings(markdown_text: str, levels: int = 1) -> str:
    """Demota TODOS los headings (#, ##, ###...) `levels` niveles hacia abajo."""
    out_lines: list[str] = []
    for line in markdown_text.splitlines():
        m = re.match(r"^(#{1,6})(\s)", line)
        if m:
            new_level = "#" * min(len(m.group(1)) + levels, 6)
            line = new_level + line[len(m.group(1)) :]
        out_lines.append(line)
    return "\n".join(out_lines)


def strip_next_class_section(markdown_text: str) -> str:
    """Elimina la sección '## ➡️ Siguiente clase' (y todo su cuerpo hasta el próximo H2 o EOF)."""
    lines = markdown_text.splitlines()
    out: list[str] = []
    skipping = False
    for line in lines:
        if NEXT_CLASS_HEADING_RE.match(line):
            skipping = True
            continue
        if skipping:
            # Termina al toparse con otro H2/H1.
            if re.match(r"^#{1,2}\s", line):
                skipping = False
                out.append(line)
            continue
        out.append(line)
    return "\n".join(out)


def class_anchor(class_dir: Path) -> str:
    """Anchor estable para una clase basado en su carpeta."""
    return f"clase-{class_dir.name}"


def class_number(class_dir: Path) -> str:
    """Extrae los 3 dígitos del prefijo de la carpeta de clase."""
    m = CLASS_DIR_RE.match(class_dir.name)
    return m.group(1) if m else class_dir.name[:3]


def class_title(class_dir: Path) -> str:
    """Lee el H1 del README de la clase."""
    readme = class_dir / "README.md"
    try:
        for line in readme.read_text(encoding="utf-8").splitlines():
            if line.startswith("# "):
                return clean_text(line[2:])
    except Exception:
        pass
    return class_dir.name


# ---------------------------------------------------------------------------
# PDFs
# ---------------------------------------------------------------------------


def build_part_markdown(
    part_dir: Path,
    part_title: str,
    classes: list[Path],
    *,
    base_class_heading_level: int = 2,
) -> str:
    """Arma el markdown grande de una parte para alimentar render_markdown_text.

    base_class_heading_level=2 -> cada clase es H2, sus headings internos se
    demotan 1 nivel (H1->H2 ya consumido; sus H2 pasan a H3, etc).
    base_class_heading_level=3 -> usado por el bundle unificado (parte=H2, clase=H3,
    headings internos demotados 2 niveles).
    """
    parts: list[str] = []
    parts.append(f"# {part_title}")
    parts.append("")
    parts.append(f"> {len(classes)} clases · bundle consolidado del currículo v3.")
    parts.append("")

    # TOC interno de la parte.
    parts.append("## Índice de clases")
    parts.append("")
    for class_dir in classes:
        num = class_number(class_dir)
        title = class_title(class_dir)
        parts.append(f"- Clase {num} — {title}")
    parts.append("")
    parts.append("---")
    parts.append("")

    demote_levels = base_class_heading_level - 1  # cuántos niveles se demota cada H1 hijo.

    for class_dir in classes:
        readme = class_dir / "README.md"
        try:
            content = readme.read_text(encoding="utf-8")
        except Exception:
            continue
        content = strip_next_class_section(content)
        # Quitar el H1 original — lo reemplazamos por el heading correcto.
        h1_re = re.compile(r"^#\s+.+$", re.MULTILINE)
        first_h1 = h1_re.search(content)
        if first_h1:
            content = content[: first_h1.start()] + content[first_h1.end() :]
            content = content.lstrip("\n")

        title = class_title(class_dir)
        num = class_number(class_dir)
        heading_hashes = "#" * base_class_heading_level
        parts.append("---")
        parts.append("")
        parts.append(f"{heading_hashes} Clase {num} — {title}")
        parts.append("")
        # Demotar todos los headings restantes para no chocar con jerarquía global.
        # base_class_heading_level=2: contenido empezaba en H2 (subtítulo), demoto 1 → H3.
        # base_class_heading_level=3: contenido empezaba en H2, demoto 2 → H4.
        demoted = demote_headings(content, levels=demote_levels + 1)
        parts.append(demoted.strip())
        parts.append("")

    parts.append("---")
    parts.append("")
    parts.append("## Cierre de la parte")
    parts.append("")
    parts.append(
        f"Fin del bundle consolidado de **{part_title}** · {len(classes)} clases."
    )
    parts.append("")
    return "\n".join(parts).rstrip() + "\n"


def build_part_pdf(part_dir: Path, output_path: Path) -> None:
    """PDF consolidado de una parte."""
    title, _subtitle = read_part_meta(part_dir)
    classes = classes_in_part(part_dir)
    m = PARTE_DIR_RE.match(part_dir.name)
    parte_num = int(m.group(1)) if m else 0
    markdown = build_part_markdown(part_dir, title, classes, base_class_heading_level=2)
    render_markdown_text(
        markdown_text=markdown,
        output_path=output_path,
        title=title,
        subtitle=f"{len(classes)} clases · Parte {parte_num} del programa",
        style="print",
    )


def build_unified_pdf(output_path: Path) -> None:
    """PDF unificado del curso completo."""
    parts = list_parts()
    total_classes = sum(len(classes_in_part(p[2])) for p in parts)

    out: list[str] = []
    out.append("# Python Data Science Program — currículo completo")
    out.append("")
    out.append(
        f"> {total_classes} clases · {len(parts)} partes · bundle unificado del currículo v3."
    )
    out.append("")
    out.append("## Índice general")
    out.append("")
    for parte_num, _slug, part_dir in parts:
        ptitle, _ = read_part_meta(part_dir)
        n = len(classes_in_part(part_dir))
        out.append(f"- **Parte {parte_num}** — {ptitle} ({n} clases)")
    out.append("")
    out.append("---")
    out.append("")

    for parte_num, _slug, part_dir in parts:
        ptitle, _ = read_part_meta(part_dir)
        classes = classes_in_part(part_dir)
        out.append("---")
        out.append("")
        out.append(f"## Parte {parte_num} — {ptitle}")
        out.append("")
        out.append(f"> {len(classes)} clases en esta parte.")
        out.append("")
        # Reutilizamos lógica de build_part_markdown pero con headings desplazados:
        # cada clase pasa a H3 (los del README se demotan 2 niveles).
        for class_dir in classes:
            readme = class_dir / "README.md"
            try:
                content = readme.read_text(encoding="utf-8")
            except Exception:
                continue
            content = strip_next_class_section(content)
            h1_re = re.compile(r"^#\s+.+$", re.MULTILINE)
            first_h1 = h1_re.search(content)
            if first_h1:
                content = content[: first_h1.start()] + content[first_h1.end() :]
                content = content.lstrip("\n")
            num = class_number(class_dir)
            title = class_title(class_dir)
            out.append(f"### Clase {num} — {title}")
            out.append("")
            out.append(demote_headings(content, levels=3).strip())
            out.append("")

    out.append("---")
    out.append("")
    out.append("## Fin del programa")
    out.append("")
    out.append(
        f"Bundle unificado generado desde {total_classes} clases en {len(parts)} partes."
    )
    out.append("")

    render_markdown_text(
        markdown_text="\n".join(out).rstrip() + "\n",
        output_path=output_path,
        title="Python Data Science Program — currículo completo",
        subtitle=f"{total_classes} clases · {len(parts)} partes · v3.7.0",
        style="print",
    )


# ---------------------------------------------------------------------------
# PPTX
# ---------------------------------------------------------------------------


def _blank_layout(prs: Presentation):
    return prs.slide_layouts[6]


def add_cover_slide(prs: Presentation, title: str, subtitle: str) -> None:
    """Slide de portada grande (centrado)."""
    slide = prs.slides.add_slide(_blank_layout(prs))
    add_background(slide)
    add_panel(slide, Inches(0.8), Inches(1.95), Inches(11.7), Inches(4.9))

    title_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.5), Inches(11.3), Inches(2.0))
    frame = title_box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = frame.paragraphs[0]
    p.text = title
    style_paragraph(
        p, font_name="Segoe UI Semibold", font_size=36, color=COLOR_TEXT, bold=True
    )

    if subtitle:
        sub_box = slide.shapes.add_textbox(
            Inches(1.0), Inches(4.7), Inches(11.3), Inches(1.2)
        )
        sframe = sub_box.text_frame
        sframe.clear()
        sframe.word_wrap = True
        sp = sframe.paragraphs[0]
        sp.text = subtitle
        style_paragraph(sp, font_name="Segoe UI", font_size=18, color=COLOR_MUTED)


def add_toc_slide(prs: Presentation, title: str, items: list[str]) -> None:
    """Slide con lista (puede paginarse en 2 columnas si hay demasiados items)."""
    # Si más de 14, partimos en chunks de 14.
    chunks: list[list[str]]
    if len(items) <= 14:
        chunks = [items]
    else:
        chunks = [items[i : i + 14] for i in range(0, len(items), 14)]

    for idx, chunk in enumerate(chunks):
        slide = prs.slides.add_slide(_blank_layout(prs))
        subtitle = ""
        if len(chunks) > 1:
            subtitle = f"Página {idx + 1} de {len(chunks)}"
        add_bullet_slide(slide, title, chunk, subtitle)


def add_divider_slide(prs: Presentation, label: str, sublabel: str = "") -> None:
    """Slide divisor con fondo sólido oscuro (color de la barra superior)."""
    slide = prs.slides.add_slide(_blank_layout(prs))
    # Fondo sólido completo.
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_HEADER_BAR

    # Banda de acento horizontal.
    accent = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, Inches(3.6), Inches(13.333), Inches(0.12)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = COLOR_ACCENT
    accent.line.color.rgb = COLOR_ACCENT

    # Título centrado.
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(2.4), Inches(12.3), Inches(1.2)
    )
    frame = title_box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = frame.paragraphs[0]
    p.text = label
    p.alignment = 2  # CENTER
    style_paragraph(
        p, font_name="Segoe UI Semibold", font_size=40, color=COLOR_BG, bold=True
    )

    if sublabel:
        sub_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(3.85), Inches(12.3), Inches(0.8)
        )
        sframe = sub_box.text_frame
        sframe.clear()
        sframe.word_wrap = True
        sp = sframe.paragraphs[0]
        sp.text = sublabel
        sp.alignment = 2
        style_paragraph(sp, font_name="Segoe UI", font_size=18, color=COLOR_PANEL)


def add_class_slides(prs: Presentation, class_dir: Path) -> None:
    """Agrega los 4 slides reducidos por clase: objetivo/resultados, temas, código, ejercicios/homework."""
    readme = class_dir / "README.md"
    if not readme.exists():
        return
    try:
        readme_data = parse_class_readme(readme)
    except Exception:
        return
    nb_data = parse_notebook(class_dir / "notebook.ipynb")

    title = readme_data["title"]
    subtitle = readme_data["subtitle"]

    objective = get_section(readme_data, "Objetivo")
    outcomes_body = get_section(readme_data, "Resultados de aprendizaje", "Resultados")
    outcomes = list_items(outcomes_body)
    topics_body = get_section(readme_data, "Temas")
    topics = topics_from_section(topics_body)
    exercises_body = get_section(readme_data, "Ejercicios")
    exercises = exercises_from_section(exercises_body)
    homework_body = get_section(readme_data, "Homework")
    homework = list_items(homework_body)
    if not homework and homework_body:
        first_para = next(
            (p.strip() for p in homework_body.split("\n\n") if p.strip()), ""
        )
        if first_para:
            homework = [clean_text(first_para)]

    blank = _blank_layout(prs)

    # 1) Objetivo + Resultados.
    slide = prs.slides.add_slide(blank)
    add_two_column_slide(
        slide,
        f"{title} — Objetivo · Resultados",
        "Objetivo",
        [objective.strip().split("\n\n")[0]] if objective else [],
        "Resultados de aprendizaje",
        outcomes[:5],
        subtitle,
    )

    # 2) Temas.
    slide = prs.slides.add_slide(blank)
    add_bullet_slide(slide, f"{title} — Temas", topics[:7], "Recorrido de la sesión.")

    # 3) Código de muestra (solo si hay).
    if nb_data.get("first_code"):
        slide = prs.slides.add_slide(blank)
        intro = nb_data.get("first_code_intro") or "Primer bloque ejecutable del notebook."
        add_code_slide(slide, f"{title} — Código", intro, nb_data["first_code"])

    # 4) Ejercicios + Homework.
    slide = prs.slides.add_slide(blank)
    add_two_column_slide(
        slide,
        f"{title} — Ejercicios · Homework",
        "Ejercicios",
        exercises[:5],
        "Homework verificable",
        homework[:5],
        "Práctica guiada + entrega.",
    )


def add_part_into_presentation(
    prs: Presentation,
    part_dir: Path,
    parte_num: int,
    *,
    include_part_cover: bool = True,
    include_part_closing: bool = True,
    closing_text: str | None = None,
) -> int:
    """Agrega los slides de una parte completa a una presentación existente."""
    title, _ = read_part_meta(part_dir)
    classes = classes_in_part(part_dir)

    if include_part_cover:
        add_cover_slide(
            prs,
            f"Parte {parte_num} — {title}",
            f"{len(classes)} clases · bundle consolidado",
        )
        toc_items = [f"Clase {class_number(c)} — {class_title(c)}" for c in classes]
        add_toc_slide(prs, f"Parte {parte_num} — Índice", toc_items)

    for class_dir in classes:
        num = class_number(class_dir)
        ctitle = class_title(class_dir)
        add_divider_slide(prs, f"Clase {num}", ctitle)
        add_class_slides(prs, class_dir)

    if include_part_closing:
        text = closing_text or f"Fin de la Parte {parte_num} — siguiente: Parte {parte_num + 1}"
        slide = prs.slides.add_slide(_blank_layout(prs))
        add_title(slide, f"Fin de la Parte {parte_num}", title)
        add_panel(slide, Inches(0.8), Inches(1.95), Inches(11.7), Inches(4.9))
        box = slide.shapes.add_textbox(Inches(1.0), Inches(3.0), Inches(11.3), Inches(2.5))
        frame = box.text_frame
        frame.clear()
        frame.word_wrap = True
        p = frame.paragraphs[0]
        p.text = text
        style_paragraph(
            p, font_name="Segoe UI Semibold", font_size=26, color=COLOR_TEXT, bold=True
        )

    return len(classes)


def _new_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs


def build_part_pptx(part_dir: Path, output_path: Path) -> None:
    """PPTX consolidado de una parte."""
    m = PARTE_DIR_RE.match(part_dir.name)
    parte_num = int(m.group(1)) if m else 0
    title, subtitle = read_part_meta(part_dir)
    classes = classes_in_part(part_dir)

    prs = _new_presentation()

    # Portada propia del bundle de parte.
    add_cover_slide(
        prs,
        f"Parte {parte_num} — {title}",
        subtitle or f"{len(classes)} clases · bundle consolidado",
    )
    # TOC.
    toc_items = [f"Clase {class_number(c)} — {class_title(c)}" for c in classes]
    add_toc_slide(prs, f"Parte {parte_num} — Índice", toc_items)

    # Determinar texto de cierre.
    all_parts = list_parts()
    max_part = max(p[0] for p in all_parts) if all_parts else parte_num
    if parte_num >= max_part:
        closing = "Fin del programa — cierre del currículo completo."
    else:
        closing = f"Fin de la Parte {parte_num} — siguiente: Parte {parte_num + 1}"

    # Clases (sin re-añadir portada/TOC, ya están).
    for class_dir in classes:
        num = class_number(class_dir)
        ctitle = class_title(class_dir)
        add_divider_slide(prs, f"Clase {num}", ctitle)
        add_class_slides(prs, class_dir)

    # Cierre.
    slide = prs.slides.add_slide(_blank_layout(prs))
    add_title(slide, f"Fin de la Parte {parte_num}", title)
    add_panel(slide, Inches(0.8), Inches(1.95), Inches(11.7), Inches(4.9))
    box = slide.shapes.add_textbox(Inches(1.0), Inches(3.0), Inches(11.3), Inches(2.5))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    p = frame.paragraphs[0]
    p.text = closing
    style_paragraph(
        p, font_name="Segoe UI Semibold", font_size=26, color=COLOR_TEXT, bold=True
    )

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)


def build_unified_pptx(output_path: Path) -> None:
    """PPTX unificado: una sola presentación con todas las partes."""
    parts = list_parts()
    total_classes = sum(len(classes_in_part(p[2])) for p in parts)

    prs = _new_presentation()
    add_cover_slide(
        prs,
        "Python Data Science Program — currículo completo",
        f"{total_classes} clases · {len(parts)} partes · v3.7.0",
    )

    # TOC general (partes).
    toc_items = []
    for parte_num, _slug, part_dir in parts:
        ptitle, _ = read_part_meta(part_dir)
        n = len(classes_in_part(part_dir))
        toc_items.append(f"Parte {parte_num} — {ptitle} ({n} clases)")
    add_toc_slide(prs, "Índice general del programa", toc_items)

    # Cada parte, sin portada individual ni cierre individual para no
    # inflar el deck con slides repetidas — pero sí dejamos divisor de parte.
    max_part = max(p[0] for p in parts) if parts else 0
    for parte_num, _slug, part_dir in parts:
        ptitle, _ = read_part_meta(part_dir)
        n = len(classes_in_part(part_dir))
        add_divider_slide(prs, f"Parte {parte_num}", f"{ptitle} · {n} clases")
        add_part_into_presentation(
            prs,
            part_dir,
            parte_num,
            include_part_cover=False,
            include_part_closing=False,
        )

    # Cierre del programa.
    add_divider_slide(
        prs,
        "Fin del programa",
        f"{total_classes} clases · {len(parts)} partes · gracias",
    )

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)


# ---------------------------------------------------------------------------
# Orquestación
# ---------------------------------------------------------------------------


def _report_ok(path: Path) -> None:
    try:
        size_kb = path.stat().st_size / 1024
    except OSError:
        size_kb = 0
    print(f"[OK] {path.relative_to(BASE_DIR)} ({size_kb:,.1f} KB)")


def _report_err(label: str, exc: Exception) -> None:
    print(
        f"[ERROR] {label}: {exc.__class__.__name__}: {exc}",
        file=sys.stderr,
    )
    traceback.print_exc(file=sys.stderr)


def run(
    *,
    parts_to_build: list[tuple[int, str, Path]],
    build_unified: bool,
    skip_pdf: bool,
    skip_pptx: bool,
) -> tuple[int, int]:
    ok = 0
    errors = 0

    for parte_num, slug, part_dir in parts_to_build:
        if not skip_pdf:
            pdf_path = PDF_PARTS_DIR / f"parte-{parte_num}-{slug}-completa.pdf"
            try:
                build_part_pdf(part_dir, pdf_path)
                _report_ok(pdf_path)
                ok += 1
            except Exception as exc:
                errors += 1
                _report_err(f"PDF parte {parte_num}", exc)
        if not skip_pptx:
            pptx_path = PPTX_PARTS_DIR / f"parte-{parte_num}-{slug}-completa.pptx"
            try:
                build_part_pptx(part_dir, pptx_path)
                _report_ok(pptx_path)
                ok += 1
            except Exception as exc:
                errors += 1
                _report_err(f"PPTX parte {parte_num}", exc)

    if build_unified:
        if not skip_pdf:
            try:
                build_unified_pdf(PDF_UNIFIED)
                _report_ok(PDF_UNIFIED)
                ok += 1
            except Exception as exc:
                errors += 1
                _report_err("PDF unificado", exc)
        if not skip_pptx:
            try:
                build_unified_pptx(PPTX_UNIFIED)
                _report_ok(PPTX_UNIFIED)
                ok += 1
            except Exception as exc:
                errors += 1
                _report_err("PPTX unificado", exc)

    return ok, errors


def main() -> None:
    parser = argparse.ArgumentParser(
        description="Genera PDFs y PPTX consolidados por parte y unificado del curso."
    )
    parser.add_argument("--parte", type=int, default=None, help="Solo esta parte (0-8).")
    parser.add_argument("--unified", action="store_true", help="Solo el bundle unificado.")
    parser.add_argument("--all", action="store_true", help="Todos los bundles + unificado (default).")
    parser.add_argument("--skip-pdf", action="store_true", help="No generar PDFs.")
    parser.add_argument("--skip-pptx", action="store_true", help="No generar PPTX.")
    args = parser.parse_args()

    all_parts = list_parts()
    if not all_parts:
        raise SystemExit(f"No se encontraron partes en {CLASSES_DIR}.")

    # Resolver qué se construye.
    if args.parte is not None and not args.unified:
        parts_to_build = [p for p in all_parts if p[0] == args.parte]
        build_unified = False
        if not parts_to_build:
            raise SystemExit(f"No existe parte {args.parte}.")
    elif args.unified and args.parte is None:
        parts_to_build = []
        build_unified = True
    elif args.unified and args.parte is not None:
        parts_to_build = [p for p in all_parts if p[0] == args.parte]
        build_unified = True
    else:
        # --all o default.
        parts_to_build = all_parts
        build_unified = True

    n_parts = len(parts_to_build)
    print(
        f"Construyendo bundles: {n_parts} parte(s) "
        f"{'+ unificado' if build_unified else ''} · "
        f"{'sin PDF ' if args.skip_pdf else ''}"
        f"{'sin PPTX ' if args.skip_pptx else ''}"
    )

    ok, errors = run(
        parts_to_build=parts_to_build,
        build_unified=build_unified,
        skip_pdf=args.skip_pdf,
        skip_pptx=args.skip_pptx,
    )

    print(f"\nResumen: {ok} OK · {errors} errores.")
    if errors:
        sys.exit(1)


if __name__ == "__main__":
    main()
