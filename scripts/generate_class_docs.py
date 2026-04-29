"""
Genera PDFs (guia-explicativa) y PPTXs (presentacion) para las clases 13-30
del bootcamp Python Data Science.
"""
import os
import re
from pathlib import Path

BASE = Path(__file__).parent.parent
CLASSES_DIR = BASE / "classes"
PDF_OUT_DIR = BASE / "docs" / "pdfs" / "classes"
PPTX_OUT_DIR = BASE / "docs" / "presentaciones" / "classes"

PDF_OUT_DIR.mkdir(parents=True, exist_ok=True)
PPTX_OUT_DIR.mkdir(parents=True, exist_ok=True)

# Classes to generate (13-30)
NEW_CLASSES = [d for d in sorted(os.listdir(CLASSES_DIR))
               if os.path.isdir(CLASSES_DIR / d)
               and d[:2].isdigit() and int(d[:2]) >= 13]

def clean_md(text):
    """Strip markdown syntax for plain text rendering."""
    text = re.sub(r'^#{1,6}\s+', '', text, flags=re.MULTILINE)
    text = re.sub(r'\*\*(.+?)\*\*', r'\1', text)
    text = re.sub(r'\*(.+?)\*', r'\1', text)
    text = re.sub(r'`{1,3}[^`]*`{1,3}', '', text)
    text = re.sub(r'^\s*[-*|>]\s*', '', text, flags=re.MULTILINE)
    text = re.sub(r'\[([^\]]+)\]\([^)]+\)', r'\1', text)
    text = re.sub(r'\n{3,}', '\n\n', text)
    return text.strip()

def read_file(path):
    if path.exists():
        return path.read_text(encoding='utf-8', errors='replace')
    return ''

# ─── PDF GENERATION ──────────────────────────────────────────────────────────
try:
    from fpdf import FPDF

    class BootcampPDF(FPDF):
        def header(self):
            self.set_font('Helvetica', 'B', 10)
            self.set_text_color(15, 61, 62)
            self.cell(0, 8, 'Bootcamp Python para Data Science', align='R')
            self.ln(12)

        def footer(self):
            self.set_y(-15)
            self.set_font('Helvetica', '', 8)
            self.set_text_color(128)
            self.cell(0, 10, f'Pagina {self.page_no()}', align='C')

        def section_title(self, title):
            self.set_font('Helvetica', 'B', 13)
            self.set_text_color(15, 61, 62)
            self.set_fill_color(230, 245, 245)
            self.cell(0, 9, title, fill=True, ln=True)
            self.ln(3)

        def body_text(self, text, size=10):
            self.set_font('Helvetica', '', size)
            self.set_text_color(30, 30, 30)
            lines = text.split('\n')
            for line in lines:
                line = line.strip()
                if not line:
                    self.ln(3)
                    continue
                # Encode for latin-1 (fpdf2 default), replace unmappable chars
                safe = line.encode('latin-1', 'replace').decode('latin-1')
                self.multi_cell(0, 6, safe)

    def generate_pdf(cls_dir, cls_name):
        readme  = clean_md(read_file(cls_dir / 'README.md'))
        teoria  = clean_md(read_file(cls_dir / 'teoria.md'))
        ejercs  = clean_md(read_file(cls_dir / 'ejercicios.md'))
        hw      = clean_md(read_file(cls_dir / 'homework.md'))
        pregs   = clean_md(read_file(cls_dir / 'preguntas.md'))
        tecno   = clean_md(read_file(cls_dir / 'tecnologias.md'))

        # Extract class number and title from folder name
        num = cls_name[:2]
        slug = cls_name[3:] if len(cls_name) > 2 else cls_name

        pdf = BootcampPDF()
        pdf.set_auto_page_break(auto=True, margin=15)
        pdf.add_page()

        # Title
        pdf.set_font('Helvetica', 'B', 18)
        pdf.set_text_color(15, 61, 62)
        title_safe = f'Clase {num} - Guia Explicativa'.encode('latin-1', 'replace').decode('latin-1')
        pdf.multi_cell(0, 10, title_safe, align='C')
        pdf.set_font('Helvetica', '', 12)
        pdf.set_text_color(80, 80, 80)
        slug_safe = slug.replace('-', ' ').title().encode('latin-1', 'replace').decode('latin-1')
        pdf.cell(0, 8, slug_safe[:80], ln=True, align='C')
        pdf.ln(8)

        # Sections
        if readme:
            pdf.section_title('Ficha de la clase')
            pdf.body_text(readme[:1200])

        if teoria:
            pdf.section_title('Base teorica')
            pdf.body_text(teoria[:2500])

        if ejercs:
            pdf.section_title('Ejercicios')
            pdf.body_text(ejercs[:1200])

        if hw:
            pdf.section_title('Tarea')
            pdf.body_text(hw[:800])

        if pregs:
            pdf.section_title('Preguntas de evaluacion')
            pdf.body_text(pregs[:1500])

        if tecno:
            pdf.section_title('Tecnologias complementarias')
            pdf.body_text(tecno[:1000])

        # Save
        fname = f'clase-{cls_name}-guia-explicativa.pdf'
        out1 = cls_dir / fname
        out2 = PDF_OUT_DIR / fname
        pdf.output(str(out1))
        pdf.output(str(out2))
        print(f'  PDF: {fname}')

    print('Generating PDFs...')
    for cls in NEW_CLASSES:
        generate_pdf(CLASSES_DIR / cls, cls)
    print(f'PDFs done: {len(NEW_CLASSES)} files')

except Exception as e:
    print(f'PDF generation error: {e}')
    import traceback; traceback.print_exc()


# ─── PPTX GENERATION ─────────────────────────────────────────────────────────
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    TEAL      = RGBColor(0x0F, 0x3D, 0x3E)
    TEAL_LIGHT= RGBColor(0xE6, 0xF5, 0xF5)
    WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
    GRAY      = RGBColor(0x50, 0x50, 0x50)

    def add_slide(prs, layout_idx=6):
        layout = prs.slide_layouts[layout_idx]
        return prs.slides.add_slide(layout)

    def set_shape_text(shape, text, size=18, bold=False, color=None):
        tf = shape.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        if color:
            run.font.color.rgb = color

    def fill_shape(shape, color):
        shape.fill.solid()
        shape.fill.fore_color.rgb = color

    def parse_slides_md(text):
        """Extract section titles and bullet points from slides.md."""
        sections = []
        current = None
        for line in text.split('\n'):
            line = line.strip()
            if line.startswith('## ') or line.startswith('# '):
                if current:
                    sections.append(current)
                current = {'title': line.lstrip('#').strip(), 'bullets': []}
            elif line.startswith('- ') or line.startswith('* '):
                if current:
                    current['bullets'].append(line[2:].strip())
            elif line.startswith('|') and current:
                pass  # skip tables
        if current:
            sections.append(current)
        return sections

    def generate_pptx(cls_dir, cls_name):
        slides_text = read_file(cls_dir / 'slides.md')
        readme_text = read_file(cls_dir / 'README.md')

        # Extract objective and key topics from README
        objective = ''
        topics = []
        for line in readme_text.split('\n'):
            if '## 🎯 Objetivo' in line or '## Objetivo' in line:
                objective = ''
            if objective == '' and line.strip() and not line.startswith('#') and not line.startswith('-'):
                if objective == '':
                    objective = line.strip()
            if line.strip().startswith('- ') and '## 🧭' not in line:
                topics.append(line.strip()[2:])

        sections = parse_slides_md(slides_text)

        num   = cls_name[:2]
        slug  = cls_name[3:].replace('-', ' ').title() if len(cls_name) > 2 else cls_name

        prs = Presentation()
        prs.slide_width  = Inches(13.33)
        prs.slide_height = Inches(7.5)

        # ── Slide 1: Title ────────────────────────────────────────────────
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        bg = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
        fill_shape(bg, TEAL)

        tb = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11), Inches(1.5))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = f'Clase {num}'
        run.font.size = Pt(20)
        run.font.bold = True
        run.font.color.rgb = TEAL_LIGHT

        tb2 = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12), Inches(2.5))
        tf2 = tb2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        run2 = p2.add_run()
        run2.text = slug
        run2.font.size = Pt(32)
        run2.font.bold = True
        run2.font.color.rgb = WHITE

        tb3 = slide.shapes.add_textbox(Inches(1), Inches(5.2), Inches(11), Inches(1))
        tf3 = tb3.text_frame
        p3 = tf3.paragraphs[0]
        run3 = p3.add_run()
        run3.text = 'Bootcamp Python para Data Science'
        run3.font.size = Pt(14)
        run3.font.color.rgb = TEAL_LIGHT

        # ── Slide 2: Objetivo ────────────────────────────────────────────
        slide2 = prs.slides.add_slide(prs.slide_layouts[6])
        hdr = slide2.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(1.3))
        fill_shape(hdr, TEAL)
        tb_h = slide2.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(1))
        set_shape_text(tb_h, 'Objetivo de la clase', 24, True, WHITE)

        tb_o = slide2.shapes.add_textbox(Inches(1), Inches(1.8), Inches(11), Inches(3))
        tf_o = tb_o.text_frame
        tf_o.word_wrap = True
        p_o = tf_o.paragraphs[0]
        run_o = p_o.add_run()
        run_o.text = objective[:300] if objective else 'Ver README.md para el objetivo de esta clase.'
        run_o.font.size = Pt(18)
        run_o.font.color.rgb = TEAL

        # ── Slide 3: Temas clave ─────────────────────────────────────────
        if topics:
            slide3 = prs.slides.add_slide(prs.slide_layouts[6])
            hdr3 = slide3.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(1.3))
            fill_shape(hdr3, TEAL)
            tb_h3 = slide3.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(1))
            set_shape_text(tb_h3, 'Temas clave', 24, True, WHITE)

            tb_t = slide3.shapes.add_textbox(Inches(1), Inches(1.8), Inches(11), Inches(4.5))
            tf_t = tb_t.text_frame
            tf_t.word_wrap = True
            for i, topic in enumerate(topics[:8]):
                p_t = tf_t.paragraphs[0] if i == 0 else tf_t.add_paragraph()
                run_t = p_t.add_run()
                run_t.text = f'• {topic}'
                run_t.font.size = Pt(16)
                run_t.font.color.rgb = GRAY

        # ── Content slides from slides.md ────────────────────────────────
        for sec in sections[:8]:
            if not sec['title'] or sec['title'].lower() in ('apertura', 'cierre esperado'):
                continue
            sl = prs.slides.add_slide(prs.slide_layouts[6])
            hdr_s = sl.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(1.3))
            fill_shape(hdr_s, TEAL)
            tb_hs = sl.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(1))
            set_shape_text(tb_hs, sec['title'][:60], 22, True, WHITE)

            if sec['bullets']:
                tb_b = sl.shapes.add_textbox(Inches(1), Inches(1.8), Inches(11), Inches(4.5))
                tf_b = tb_b.text_frame
                tf_b.word_wrap = True
                for i, bullet in enumerate(sec['bullets'][:7]):
                    pb = tf_b.paragraphs[0] if i == 0 else tf_b.add_paragraph()
                    rb = pb.add_run()
                    rb.text = f'• {bullet}'
                    rb.font.size = Pt(15)
                    rb.font.color.rgb = GRAY

        # ── Final slide ──────────────────────────────────────────────────
        sl_f = prs.slides.add_slide(prs.slide_layouts[6])
        bg_f = sl_f.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
        fill_shape(bg_f, TEAL)
        tb_f = sl_f.shapes.add_textbox(Inches(2), Inches(2.5), Inches(9), Inches(2))
        set_shape_text(tb_f, 'A practicar', 40, True, WHITE)
        tb_f2 = sl_f.shapes.add_textbox(Inches(2), Inches(4.2), Inches(9), Inches(1))
        set_shape_text(tb_f2, f'Clase {num} — {slug}', 18, False, TEAL_LIGHT)

        fname = f'clase-{cls_name}-presentacion.pptx'
        out1 = cls_dir / fname
        out2 = PPTX_OUT_DIR / fname
        prs.save(str(out1))
        prs.save(str(out2))
        print(f'  PPTX: {fname}')

    print('\nGenerating PPTXs...')
    for cls in NEW_CLASSES:
        generate_pptx(CLASSES_DIR / cls, cls)
    print(f'PPTXs done: {len(NEW_CLASSES)} files')

except Exception as e:
    print(f'PPTX generation error: {e}')
    import traceback; traceback.print_exc()

print('\nAll done.')
