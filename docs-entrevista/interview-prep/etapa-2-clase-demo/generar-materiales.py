#!/usr/bin/env python3
"""
Generador de materiales para Etapa 2 — Presentación de Clase
Crea PPTX + PDF de preparación para la Clase 11: Evaluación y Pipelines
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
import json

# ============================================================================
# CONFIGURACIÓN DE COLORES
# ============================================================================

COLOR_PRIMARY = RGBColor(15, 61, 62)      # Teal oscuro
COLOR_ACCENT = RGBColor(255, 107, 53)    # Naranja
COLOR_TEXT = RGBColor(33, 33, 33)        # Gris oscuro
COLOR_LIGHT = RGBColor(242, 242, 242)    # Gris claro
COLOR_SUCCESS = RGBColor(76, 175, 80)    # Verde
COLOR_WARNING = RGBColor(255, 152, 0)    # Naranja claro

# ============================================================================
# ESTRUCTURA DE LA CLASE 11
# ============================================================================

CLASE_TITULO = "Evaluación y Pipelines — Optimización de Modelos"
CLASE_NUM = "11"

SLIDES = [
    {
        "titulo": "Evaluación y Pipelines",
        "subtitulo": "Clase 11 — Optimización y automatización en ML",
        "tipo": "titulo",
        "notas": "Abre con energía. Esta clase cierra el ciclo: aprendemos a evaluar correctamente y automatizar workflows."
    },
    {
        "titulo": "Agenda de hoy",
        "contenido": [
            "• Validación cruzada (cross-validation) — por qué no es suficiente train/test",
            "• Métricas avanzadas — ROC-AUC, precision-recall, F1",
            "• GridSearchCV — busca automática de hiperparámetros",
            "• Pipelines — encadena preprocesamiento + modelo",
            "• Ejercicios prácticos — optimización de un modelo real",
            "• Proyecto final — integración de todo lo aprendido"
        ],
        "tipo": "contenido",
        "notas": "Muestra esta agenda al inicio para que los alumnos sepan qué esperar. 5 min."
    },
    {
        "titulo": "El problema: overfitting",
        "contenido": [
            "Train accuracy: 95% ❌ (muy bueno para ser verdad)",
            "Test accuracy: 62% ❌ (modelo no generaliza)",
            "",
            "¿Por qué pasó? El modelo memorizó el training set.",
            "",
            "Solución: validación cruzada (k-fold cross-validation)"
        ],
        "tipo": "contenido",
        "notas": "Dibuja curvas de learning en el pizarrón. Muestra cómo el training loss baja pero test loss sube."
    },
    {
        "titulo": "Cross-Validation (k-fold)",
        "contenido": [
            "Divide el dataset en k partes (típicamente k=5)",
            "",
            "Iterativamente:",
            "  • Usa k-1 partes para entrenar",
            "  • Usa 1 parte para validar",
            "",
            "Resultado: promedio de k scores → estimación más confiable"
        ],
        "tipo": "contenido",
        "notas": "Código de ejemplo:\nfrom sklearn.model_selection import cross_val_score\nscores = cross_val_score(model, X, y, cv=5)\nprint(f'Promedio: {scores.mean():.3f} +/- {scores.std():.3f}')"
    },
    {
        "titulo": "Código: Cross-Validation",
        "contenido": [
            "from sklearn.model_selection import cross_val_score",
            "from sklearn.tree import DecisionTreeClassifier",
            "",
            "model = DecisionTreeClassifier()",
            "scores = cross_val_score(model, X, y, cv=5)",
            "",
            "print(f'Scores: {scores}')",
            "print(f'Media: {scores.mean():.3f}')",
            "print(f'Std:   {scores.std():.3f}')"
        ],
        "tipo": "codigo",
        "notas": "Ejecuta esto en vivo en Jupyter. Muestra cada fold y el promedio."
    },
    {
        "titulo": "Métricas: confusion matrix",
        "contenido": [
            "                Predicho: No    Predicho: Sí",
            "Real: No             TN              FP",
            "Real: Sí             FN              TP",
            "",
            "TN = Verdaderos Negativos (correcto, no enfermo)",
            "TP = Verdaderos Positivos (correcto, enfermo)",
            "FP = Falsos Positivos (error, dijo enfermo pero no lo es)",
            "FN = Falsos Negativos (error, dijo no enfermo pero sí lo es)"
        ],
        "tipo": "contenido",
        "notas": "Usa un ejemplo médico para que sea intuitive. FN es peligroso (paciente enfermo sin detectar)."
    },
    {
        "titulo": "Métricas: Precision, Recall, F1",
        "contenido": [
            "Precision = TP / (TP + FP)  → de los predichos Sí, cuántos eran correctos",
            "Recall    = TP / (TP + FN)  → de los reales Sí, cuántos atrapamos",
            "F1        = 2 × (P × R) / (P + R)  → balance entre P y R",
            "",
            "Usa Precision cuando: falsos positivos son caros",
            "Usa Recall cuando: falsos negativos son caros",
            "Usa F1 cuando: ambos errores importan igual"
        ],
        "tipo": "contenido",
        "notas": "Ejemplo: spam detection (falso positivo = email importante perdido). Vs diagnóstico (falso negativo = enfermedad no detectada)."
    },
    {
        "titulo": "Código: classification_report",
        "contenido": [
            "from sklearn.metrics import classification_report",
            "",
            "y_pred = model.predict(X_test)",
            "print(classification_report(y_test, y_pred))",
            "",
            "Resultado:",
            "           precision  recall  f1-score",
            "     No       0.92     0.88      0.90",
            "     Sí       0.85     0.91      0.88",
            "  accuracy            0.89"
        ],
        "tipo": "codigo",
        "notas": "Explica qué significa cada columna. El accuracy general es menos informativo que los por-clase."
    },
    {
        "titulo": "GridSearchCV — búsqueda automática",
        "contenido": [
            "¿Qué hiperparámetro es mejor?",
            "max_depth = 3, 5, 7, 10?",
            "min_samples_split = 2, 5, 10?",
            "",
            "GridSearchCV prueba TODAS las combinaciones automáticamente",
            "y te retorna la mejor."
        ],
        "tipo": "contenido",
        "notas": "Analógico: en lugar de probar manualmente, usamos una herramienta que lo hace por nosotros."
    },
    {
        "titulo": "Código: GridSearchCV",
        "contenido": [
            "from sklearn.model_selection import GridSearchCV",
            "from sklearn.tree import DecisionTreeClassifier",
            "",
            "param_grid = {",
            "    'max_depth': [3, 5, 7, 10],",
            "    'min_samples_split': [2, 5, 10]",
            "}",
            "grid = GridSearchCV(DecisionTreeClassifier(), param_grid, cv=5)",
            "grid.fit(X_train, y_train)",
            "print(f'Mejor: {grid.best_params_}')",
            "print(f'Score: {grid.best_score_:.3f}')"
        ],
        "tipo": "codigo",
        "notas": "Esto puede tardar. Explica que GridSearchCV prueba 4 × 3 × 5 = 60 configuraciones."
    },
    {
        "titulo": "Pipelines — encadena pasos",
        "contenido": [
            "Problema sin Pipeline:",
            "  1. StandardScaler → X_train, X_test",
            "  2. PCA → X_train_scaled, X_test_scaled",
            "  3. LogisticRegression → fit/predict",
            "  (¡fácil olvidar un paso o hacerlo mal!)",
            "",
            "Solución: Pipeline",
            "  pipeline = Pipeline([",
            "      ('scaler', StandardScaler()),",
            "      ('pca', PCA()),",
            "      ('clf', LogisticRegression())",
            "  ])",
            "  pipeline.fit(X_train, y_train)  ← ¡todo en una línea!"
        ],
        "tipo": "contenido",
        "notas": "Ventajas: código limpio, menos errores, fácil de reproducir."
    },
    {
        "titulo": "Código: Pipeline + GridSearchCV",
        "contenido": [
            "from sklearn.pipeline import Pipeline",
            "from sklearn.preprocessing import StandardScaler",
            "from sklearn.decomposition import PCA",
            "from sklearn.linear_model import LogisticRegression",
            "",
            "pipe = Pipeline([",
            "    ('scaler', StandardScaler()),",
            "    ('pca', PCA()),",
            "    ('clf', LogisticRegression())",
            "])",
            "",
            "param_grid = {'pca__n_components': [2, 5, 10]}",
            "grid = GridSearchCV(pipe, param_grid, cv=5)",
            "grid.fit(X_train, y_train)"
        ],
        "tipo": "codigo",
        "notas": "Usa 'pca__n_components' para acceder parámetros del paso 'pca' en el pipeline."
    },
    {
        "titulo": "Ejercicio en vivo: dataset Iris",
        "contenido": [
            "1. Carga iris dataset",
            "2. Split train/test (80/20)",
            "3. Crea Pipeline: Scaler → LogisticRegression",
            "4. Usa GridSearchCV para optimizar C (regularización)",
            "5. Reporta mejores parámetros y accuracy final",
            "6. Genera classification_report",
            "",
            "⏱️ 15 minutos. Yo lo haré paso a paso en vivo."
        ],
        "tipo": "contenido",
        "notas": "Abre Jupyter. Vé lentamente. Deja que los alumnos escriban con vos."
    },
    {
        "titulo": "Resumen de la clase",
        "contenido": [
            "✓ Cross-Validation → estimación confiable sin overfitting",
            "✓ Métricas correctas → Precision, Recall, F1 (no solo accuracy)",
            "✓ GridSearchCV → encuentra automáticamente los mejores hiperparámetros",
            "✓ Pipelines → encadena pasos, código limpio y reproducible",
            "",
            "Siguiente: Clase 12 integra TODO en un proyecto final."
        ],
        "tipo": "contenido",
        "notas": "Resume los puntos clave. Motiva para la clase siguiente."
    },
    {
        "titulo": "Preguntas y discusión",
        "contenido": [
            "¿Cuándo usar GridSearchCV vs búsqueda manual?",
            "¿Es Pipeline sólo para clasificación?",
            "¿Cómo sé si estoy overfitting?",
            "¿Qué pasa si GridSearchCV tarda mucho tiempo?",
            "",
            "¡Pregunten! No hay preguntas tontas."
        ],
        "tipo": "contenido",
        "notas": "Deja espacio para diálogo. Esto muestra que eres accesible como docente."
    }
]

# ============================================================================
# FUNCIONES
# ============================================================================

def agregar_slide_titulo(prs, slide):
    """Slide de portada"""
    layout = prs.slide_layouts[6]  # Blank
    slide_obj = prs.slides.add_slide(layout)

    # Fondo
    background = slide_obj.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_PRIMARY

    # Título
    title_box = slide_obj.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = slide["titulo"]
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    # Subtítulo
    subtitle_box = slide_obj.shapes.add_textbox(Inches(0.5), Inches(3.7), Inches(9), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    p = subtitle_frame.paragraphs[0]
    p.text = slide["subtitulo"]
    p.font.size = Pt(28)
    p.font.color.rgb = COLOR_ACCENT

    return slide_obj

def agregar_slide_contenido(prs, slide):
    """Slide con bullets"""
    layout = prs.slide_layouts[6]  # Blank
    slide_obj = prs.slides.add_slide(layout)

    # Fondo
    background = slide_obj.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)

    # Encabezado con fondo
    header_shape = slide_obj.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(10), Inches(0.8)
    )
    header_shape.fill.solid()
    header_shape.fill.fore_color.rgb = COLOR_PRIMARY
    header_shape.line.color.rgb = COLOR_PRIMARY

    # Título
    title_box = slide_obj.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = slide["titulo"]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    # Contenido
    content_box = slide_obj.shapes.add_textbox(Inches(0.7), Inches(1.1), Inches(8.6), Inches(5))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True

    for i, line in enumerate(slide["contenido"]):
        if i > 0:
            p = text_frame.add_paragraph()
        else:
            p = text_frame.paragraphs[0]

        p.text = line
        p.font.size = Pt(18)
        p.font.color.rgb = COLOR_TEXT
        p.space_before = Pt(3)
        p.space_after = Pt(3)

        if line.startswith("•"):
            p.level = 0

    # Notas del docente
    notes_slide = slide_obj.notes_slide
    notes_text_frame = notes_slide.notes_text_frame
    notes_text_frame.text = slide.get("notas", "")

    return slide_obj

def agregar_slide_codigo(prs, slide):
    """Slide con código"""
    layout = prs.slide_layouts[6]  # Blank
    slide_obj = prs.slides.add_slide(layout)

    # Fondo
    background = slide_obj.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)

    # Encabezado
    header_shape = slide_obj.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.8))
    header_shape.fill.solid()
    header_shape.fill.fore_color.rgb = COLOR_PRIMARY
    header_shape.line.color.rgb = COLOR_PRIMARY

    # Título
    title_box = slide_obj.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = slide["titulo"]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    # Código
    code_box = slide_obj.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(9), Inches(5.5))
    text_frame = code_box.text_frame
    text_frame.word_wrap = True

    for i, line in enumerate(slide["contenido"]):
        if i > 0:
            p = text_frame.add_paragraph()
        else:
            p = text_frame.paragraphs[0]

        p.text = line
        p.font.name = "Courier New"
        p.font.size = Pt(14)
        p.font.color.rgb = COLOR_TEXT
        p.space_before = Pt(1)
        p.space_after = Pt(1)

    # Notas
    notes_slide = slide_obj.notes_slide
    notes_text_frame = notes_slide.notes_text_frame
    notes_text_frame.text = slide.get("notas", "")

    return slide_obj

# ============================================================================
# GENERAR PRESENTACIÓN
# ============================================================================

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

for slide_data in SLIDES:
    if slide_data["tipo"] == "titulo":
        agregar_slide_titulo(prs, slide_data)
    elif slide_data["tipo"] == "codigo":
        agregar_slide_codigo(prs, slide_data)
    else:
        agregar_slide_contenido(prs, slide_data)

prs.save("Clase_11_Evaluacion_y_Pipelines.pptx")
print("[OK] PPTX generado: Clase_11_Evaluacion_y_Pipelines.pptx")
