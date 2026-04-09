#!/usr/bin/env python3
"""
Generador de 3 Presentaciones PPTX (45 minutos cada una)
Optimizadas para presentación ONLINE
Clase 6 (Visualización) | Clase 9 (ML Intro) | Clase 11 (Pipelines)
"""

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

COLOR_PRIMARY = RGBColor(15, 61, 62)
COLOR_ACCENT = RGBColor(255, 107, 53)
COLOR_TEXT = RGBColor(33, 33, 33)

# ============================================================================
# CLASE 6: VISUALIZACIÓN CON MATPLOTLIB (45 min)
# ============================================================================

CLASE_6 = {
    "titulo": "Visualización con Matplotlib",
    "numero": "06",
    "slides": [
        {
            "tipo": "titulo",
            "titulo": "Visualización con Matplotlib",
            "subtitulo": "Clase 6 — Convierte datos en gráficos claros",
            "notas": "Abre con energía. Esta clase es sobre VER los datos, no solo números."
        },
        {
            "tipo": "contenido",
            "titulo": "¿Por qué visualizar?",
            "contenido": [
                "• 1000 números en una tabla = confuso",
                "• Un gráfico = comprensión inmediata",
                "",
                "Anscombe's Quartet: 4 datasets idénticos estadísticamente",
                "pero gráficamente COMPLETAMENTE diferentes",
                "",
                "→ Siempre visualiza ANTES de reportar"
            ],
            "notas": "Muestra la imagen de Anscombe si puedes (Google Images). Impacto visual fuerte."
        },
        {
            "tipo": "contenido",
            "titulo": "Tipos de gráficos básicos",
            "contenido": [
                "• Scatter plot: relación entre 2 variables",
                "• Line plot: tendencia en el tiempo",
                "• Bar chart: comparar categorías",
                "• Histogram: distribución de datos",
                "• Box plot: detectar outliers",
                "• Heatmap: correlaciones entre variables"
            ],
            "notas": "Cada tipo tiene un uso. No uses bar chart para tendencias (usa line)."
        },
        {
            "tipo": "codigo",
            "titulo": "Primer gráfico",
            "contenido": [
                "import matplotlib.pyplot as plt",
                "import numpy as np",
                "",
                "x = np.linspace(0, 10, 100)",
                "y = np.sin(x)",
                "",
                "plt.figure(figsize=(10, 6))",
                "plt.plot(x, y, linewidth=2)",
                "plt.title('Onda Seno')",
                "plt.xlabel('X')",
                "plt.ylabel('sin(x)')",
                "plt.grid(True)",
                "plt.show()"
            ],
            "notas": "Ejecuta en vivo. Muestra cómo cambiar title, labels, grid."
        },
        {
            "tipo": "contenido",
            "titulo": "Estilo y colores",
            "contenido": [
                "• plt.style.use('seaborn'): tema predefinido",
                "• color='red': especificar color",
                "• linewidth=2: grosor de línea",
                "• marker='o': puntos en los datos",
                "• alpha=0.7: transparencia",
                "",
                "plt.savefig('grafico.png', dpi=300)"
            ],
            "notas": "Más bonito = más impacto. Los gráficos feos no se ven profesionales."
        },
        {
            "tipo": "codigo",
            "titulo": "Subplots (múltiples gráficos)",
            "contenido": [
                "fig, axes = plt.subplots(2, 2, figsize=(12, 10))",
                "",
                "axes[0, 0].plot([1, 2, 3, 4])",
                "axes[0, 0].set_title('Gráfico 1')",
                "",
                "axes[0, 1].scatter([1, 2, 3], [4, 5, 6])",
                "axes[0, 1].set_title('Gráfico 2')",
                "",
                "# ... etc para los otros 2",
                "",
                "plt.tight_layout()",
                "plt.show()"
            ],
            "notas": "Para reportes: múltiples gráficos en una figura. tight_layout() evita que se superpongan."
        },
        {
            "tipo": "contenido",
            "titulo": "Ejercicio rápido (10 min)",
            "contenido": [
                "Datos: dataset Iris",
                "Tarea: Crea 2 gráficos lado a lado",
                "  1. Scatter: Sepal Length vs Sepal Width",
                "  2. Histogram: Distribución de Petal Length",
                "",
                "Bonus: Usa colores diferentes por especie"
            ],
            "notas": "Los alumnos escriben conmigo. Lentamente. Si van rápido, expande con más detalles."
        },
        {
            "tipo": "contenido",
            "titulo": "Resumen",
            "contenido": [
                "✓ Visualización es tan importante como los datos",
                "✓ Matplotlib es la herramienta estándar",
                "✓ Tipos de gráficos para cada caso",
                "✓ Estilo importa (colores, títulos, claridad)",
                "✓ Practica: datos → gráfico → insight"
            ],
            "notas": "Cierre fuerte. La visualización es cómo comunicas resultados al mundo."
        }
    ]
}

# ============================================================================
# CLASE 9: MACHINE LEARNING INTRO (45 min)
# ============================================================================

CLASE_9 = {
    "titulo": "Machine Learning Intro",
    "numero": "09",
    "slides": [
        {
            "tipo": "titulo",
            "titulo": "Machine Learning Intro",
            "subtitulo": "Clase 9 — Las máquinas aprenden de los datos",
            "notas": "Abre con una pregunta: '¿Cómo Netflix predice qué serie te gustará?' → Machine Learning"
        },
        {
            "tipo": "contenido",
            "titulo": "¿Qué es Machine Learning?",
            "contenido": [
                "Definición: Permitir que máquinas aprendan de datos sin ser programadas explícitamente",
                "",
                "3 tipos:",
                "• Supervised: datos + respuestas correctas → predice",
                "• Unsupervised: solo datos → encuentra patrones",
                "• Reinforcement: rewards → aprende a actuar"
            ],
            "notas": "Ejemplo Supervised: email spam/no spam. Unsupervised: agrupar clientes similares."
        },
        {
            "tipo": "contenido",
            "titulo": "Flujo típico",
            "contenido": [
                "1. OBTENER datos (dataset)",
                "2. EXPLORAR (¿qué hay en los datos?)",
                "3. LIMPIAR (valores nulos, outliers)",
                "4. PREPARAR (normalizar, seleccionar features)",
                "5. ENTRENAR modelo",
                "6. EVALUAR (¿qué tan bien funciona?)",
                "7. DESPLEGAR (en producción)"
            ],
            "notas": "Este flujo se repite. No es lineal — iteras constantemente."
        },
        {
            "tipo": "codigo",
            "titulo": "Primer modelo: regresión lineal",
            "contenido": [
                "from sklearn.linear_model import LinearRegression",
                "from sklearn.model_selection import train_test_split",
                "",
                "# Datos",
                "X_train, X_test, y_train, y_test = train_test_split(X, y, test_size=0.2)",
                "",
                "# Modelo",
                "modelo = LinearRegression()",
                "modelo.fit(X_train, y_train)",
                "",
                "# Predicción",
                "y_pred = modelo.predict(X_test)",
                "score = modelo.score(X_test, y_test)  # 0.0 a 1.0"
            ],
            "notas": "Regresión lineal es simple pero poderosa. Score = R2 (qué % de varianza explica)."
        },
        {
            "tipo": "contenido",
            "titulo": "Overfitting vs Underfitting",
            "contenido": [
                "UNDERFITTING: Modelo muy simple",
                "  train score: 60%, test score: 58%",
                "  → Ambos bajos (no aprendió nada)",
                "",
                "OVERFITTING: Modelo muy complejo",
                "  train score: 95%, test score: 65%",
                "  → Memorizó training, no generaliza",
                "",
                "BALANCEADO: El objetivo",
                "  train score: 88%, test score: 85%"
            ],
            "notas": "El gap (diferencia) es clave. Un gap pequeño = buen modelo."
        },
        {
            "tipo": "contenido",
            "titulo": "Algoritmos comunes",
            "contenido": [
                "Regresión: LinearRegression, Polynomial, Ridge",
                "Clasificación: LogisticRegression, DecisionTree, RandomForest",
                "Clustering: KMeans, DBSCAN",
                "",
                "No hay algoritmo 'mejor'. Depende del problema."
            ],
            "notas": "Clase 10 entra en profundidad. Aquí solo mencionas los nombres."
        },
        {
            "tipo": "codigo",
            "titulo": "Evaluación: Accuracy",
            "contenido": [
                "from sklearn.metrics import accuracy_score",
                "",
                "y_pred = modelo.predict(X_test)",
                "",
                "accuracy = accuracy_score(y_test, y_pred)",
                "print(f'Accuracy: {accuracy:.2%}')  # 85.23%",
                "",
                "# Otros scores:",
                "modelo.score(X_test, y_test)  # Lo mismo"
            ],
            "notas": "Accuracy = % correcto. Es la métrica más simple (pero no siempre la mejor)."
        },
        {
            "tipo": "contenido",
            "titulo": "Resumen",
            "contenido": [
                "✓ ML: máquinas aprenden de datos",
                "✓ 3 tipos: Supervised, Unsupervised, Reinforcement",
                "✓ Flujo: obtener → explorar → limpiar → preparar → entrenar → evaluar → desplegar",
                "✓ Overfitting es el enemigo",
                "✓ Próxima clase: algoritmos supervisados en profundidad"
            ],
            "notas": "Motivación para Clase 10. Hoy vieron el panorama, mañana los detalles."
        }
    ]
}

# ============================================================================
# CLASE 11: EVALUACIÓN Y PIPELINES (45 min) — VERSIÓN CORTA
# ============================================================================

CLASE_11 = {
    "titulo": "Evaluación y Pipelines",
    "numero": "11",
    "slides": [
        {
            "tipo": "titulo",
            "titulo": "Evaluación y Pipelines",
            "subtitulo": "Clase 11 — Optimización y automatización",
            "notas": "Esta es la penúltima clase. Cierra el ciclo de ML."
        },
        {
            "tipo": "contenido",
            "titulo": "El problema",
            "contenido": [
                "Train accuracy: 95% ❌",
                "Test accuracy: 62% ❌",
                "",
                "¿Qué pasó? Overfitting.",
                "",
                "Solución: Cross-Validation (k-fold)"
            ],
            "notas": "Dibuja curvas en la pizarra (training vs validation)."
        },
        {
            "tipo": "codigo",
            "titulo": "Cross-Validation",
            "contenido": [
                "from sklearn.model_selection import cross_val_score",
                "",
                "modelo = DecisionTreeClassifier()",
                "scores = cross_val_score(modelo, X, y, cv=5)",
                "",
                "print(f'Scores: {scores}')",
                "print(f'Media: {scores.mean():.3f} +/- {scores.std():.3f}')"
            ],
            "notas": "cv=5 = 5-fold. Usa 4/5 para entrenar, 1/5 para validar. Repite 5 veces."
        },
        {
            "tipo": "contenido",
            "titulo": "Métricas: Precision & Recall",
            "contenido": [
                "Precision = TP / (TP + FP) = de los que predije Sí, cuántos eran correctos",
                "Recall = TP / (TP + FN) = de los Sí reales, cuántos atrapé",
                "F1 = balance entre Precision y Recall",
                "",
                "Ejemplo: Diagnóstico médico",
                "Recall > Precision (no puedo perder pacientes enfermos)"
            ],
            "notas": "La métrica correcta depende del problema. Accuracy no siempre es suficiente."
        },
        {
            "tipo": "codigo",
            "titulo": "GridSearchCV",
            "contenido": [
                "from sklearn.model_selection import GridSearchCV",
                "",
                "param_grid = {'max_depth': [3, 5, 7, 10], 'min_samples_split': [2, 5, 10]}",
                "",
                "grid = GridSearchCV(DecisionTreeClassifier(), param_grid, cv=5)",
                "grid.fit(X_train, y_train)",
                "",
                "print(f'Mejor: {grid.best_params_}')",
                "print(f'Score: {grid.best_score_:.3f}')"
            ],
            "notas": "GridSearchCV prueba automáticamente todas las combinaciones. Encontra los mejores parámetros."
        },
        {
            "tipo": "contenido",
            "titulo": "Pipelines",
            "contenido": [
                "Sin Pipeline: normalizar → PCA → modelo (3 pasos, fácil olvidar uno)",
                "",
                "Con Pipeline: todo junto y coordinado",
                "Ventajas:",
                "• Código limpio",
                "• Menos errores",
                "• Reproducible"
            ],
            "notas": "Pipeline = encadenamiento automático. Ahorra errores y tiempo."
        },
        {
            "tipo": "codigo",
            "titulo": "Pipeline + GridSearchCV",
            "contenido": [
                "from sklearn.pipeline import Pipeline",
                "",
                "pipe = Pipeline([",
                "    ('scaler', StandardScaler()),",
                "    ('clf', LogisticRegression())",
                "])",
                "",
                "param_grid = {'clf__C': [0.1, 1, 10]}",
                "grid = GridSearchCV(pipe, param_grid, cv=5)",
                "grid.fit(X_train, y_train)"
            ],
            "notas": "clf__C accede parámetros del paso 'clf' dentro del Pipeline. Clave."
        },
        {
            "tipo": "contenido",
            "titulo": "Resumen",
            "contenido": [
                "✓ Cross-Validation: estimación confiable",
                "✓ Métricas correctas: Precision, Recall, F1 (contexto importa)",
                "✓ GridSearchCV: encuentra automáticamente mejores hiperparámetros",
                "✓ Pipelines: encadena pasos, código limpio y reproducible",
                "✓ Próxima: Proyecto Final integrando TODO"
            ],
            "notas": "Esta clase cierra el ciclo. Proyecto final integra todo lo aprendido."
        }
    ]
}

# ============================================================================
# FUNCIONES
# ============================================================================

def crear_slide_titulo(prs, slide_data):
    """Crea slide de portada"""
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)

    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_PRIMARY

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = slide_data["titulo"]
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    p = subtitle_frame.paragraphs[0]
    p.text = slide_data["subtitulo"]
    p.font.size = Pt(24)
    p.font.color.rgb = COLOR_ACCENT

    notes = slide.notes_slide.notes_text_frame
    notes.text = slide_data["notas"]

    return slide

def crear_slide_contenido(prs, slide_data):
    """Crea slide con bullets"""
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)

    # Fondo blanco
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)

    # Header
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.8))
    header.fill.solid()
    header.fill.fore_color.rgb = COLOR_PRIMARY
    header.line.color.rgb = COLOR_PRIMARY

    # Título
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = slide_data["titulo"]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    # Contenido
    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.1), Inches(8.6), Inches(5.5))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True

    for i, line in enumerate(slide_data["contenido"]):
        if i > 0:
            p = text_frame.add_paragraph()
        else:
            p = text_frame.paragraphs[0]

        p.text = line
        p.font.size = Pt(16)
        p.font.color.rgb = COLOR_TEXT
        p.space_before = Pt(2)
        p.space_after = Pt(2)

    notes = slide.notes_slide.notes_text_frame
    notes.text = slide_data["notas"]

    return slide

def crear_slide_codigo(prs, slide_data):
    """Crea slide con código"""
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)

    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)

    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.8))
    header.fill.solid()
    header.fill.fore_color.rgb = COLOR_PRIMARY
    header.line.color.rgb = COLOR_PRIMARY

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = slide_data["titulo"]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    # Código
    code_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(9), Inches(5.5))
    text_frame = code_box.text_frame
    text_frame.word_wrap = True

    for i, line in enumerate(slide_data["contenido"]):
        if i > 0:
            p = text_frame.add_paragraph()
        else:
            p = text_frame.paragraphs[0]

        p.text = line
        p.font.name = "Courier New"
        p.font.size = Pt(12)
        p.font.color.rgb = COLOR_TEXT
        p.space_before = Pt(1)
        p.space_after = Pt(1)

    notes = slide.notes_slide.notes_text_frame
    notes.text = slide_data["notas"]

    return slide

def generar_pptx(clase_data, filename):
    """Genera PPTX de una clase"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    for slide_data in clase_data["slides"]:
        if slide_data["tipo"] == "titulo":
            crear_slide_titulo(prs, slide_data)
        elif slide_data["tipo"] == "codigo":
            crear_slide_codigo(prs, slide_data)
        else:
            crear_slide_contenido(prs, slide_data)

    prs.save(filename)
    print(f"[OK] {filename}")

# ============================================================================
# MAIN
# ============================================================================

generar_pptx(CLASE_6, "Clase_06_Visualizacion_45min.pptx")
generar_pptx(CLASE_9, "Clase_09_ML_Intro_45min.pptx")
generar_pptx(CLASE_11, "Clase_11_Evaluacion_Pipelines_45min.pptx")

print("\n[OK] 3 presentaciones generadas (45 minutos cada una)")
